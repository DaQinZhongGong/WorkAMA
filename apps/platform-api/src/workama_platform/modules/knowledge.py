from __future__ import annotations

import asyncio
import hashlib
import io
import json
import math
import mimetypes
import re
import struct
import zipfile
from collections.abc import Iterable
from datetime import UTC, datetime
from pathlib import Path
from typing import Annotated, Any
from urllib.parse import urlsplit

import httpx
from bs4 import BeautifulSoup
from docx import Document as DocxDocument
from fastapi import APIRouter, Depends, Header, HTTPException, Request, Response
from openpyxl import load_workbook
from PIL import Image
from pydantic import BaseModel, Field, field_validator
from pypdf import PdfReader
from pptx import Presentation
from pytesseract import image_to_string

from workama_platform.core import (
    Actor,
    capability_allows,
    get_actor,
    json_dumps,
    new_id,
    pool,
    settings,
)
from workama_platform.modules.jobs import (
    ClaimedJob,
    IdempotencyConflict,
    request_cancellation,
    submit_operation,
)
from workama_platform.modules.security.service import validate_resolved_outbound_url
from workama_platform.object_store import delete_object, get_object, put_object


router = APIRouter(prefix="/api/v1", tags=["knowledge"])

DATASET_BUCKET = "workama-datasets"
MAX_DOCUMENT_BYTES = 100 * 1024 * 1024
MAX_EXTRACTED_CHARS = 10_000_000
DEFAULT_RETRIEVAL_CONFIG = {
    "top_k": 5,
    "candidate_k": 20,
    "rrf_k": 60,
    "score_threshold": 0.0,
}
SUPPORTED_SUFFIXES = {
    ".txt": "text/plain",
    ".md": "text/markdown",
    ".markdown": "text/markdown",
    ".html": "text/html",
    ".htm": "text/html",
    ".json": "application/json",
    ".csv": "text/csv",
    ".pdf": "application/pdf",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ".png": "image/png",
    ".jpg": "image/jpeg",
    ".jpeg": "image/jpeg",
    ".webp": "image/webp",
}
TEXT_MIME_TYPES = {
    "text/plain",
    "text/markdown",
    "text/html",
    "application/json",
    "text/csv",
    "application/csv",
    "application/xml",
    "text/xml",
}
IMAGE_MIME_TYPES = {"image/png", "image/jpeg", "image/webp"}


class RagProcessingError(RuntimeError):
    code = "E03001"


class RagJobCancelled(RuntimeError):
    pass


class DatasetCreate(BaseModel):
    name: str = Field(min_length=1, max_length=120)
    description: str = Field(default="", max_length=2_000)
    embedding_model: str = Field(default="workama-embed", min_length=1, max_length=120)
    retrieval_config: dict[str, Any] = Field(default_factory=dict)

    @field_validator("name")
    @classmethod
    def normalize_name(cls, value: str) -> str:
        normalized = " ".join(value.split())
        if not normalized:
            raise ValueError("name must not be empty")
        return normalized


class DatasetPatch(BaseModel):
    name: str | None = Field(default=None, min_length=1, max_length=120)
    description: str | None = Field(default=None, max_length=2_000)


class RetrievalConfigUpsert(BaseModel):
    top_k: int = Field(default=5, ge=1, le=50)
    candidate_k: int = Field(default=20, ge=5, le=200)
    rrf_k: int = Field(default=60, ge=1, le=200)
    score_threshold: float = Field(default=0.0, ge=0.0, le=1.0)

    def normalized(self) -> dict[str, Any]:
        if self.candidate_k < self.top_k:
            raise HTTPException(status_code=422, detail="candidate_k must be at least top_k")
        return self.model_dump()


class DeleteReason(BaseModel):
    reason: str = Field(min_length=3, max_length=500)


class RestoreRequest(BaseModel):
    reason: str = Field(default="User requested restore", min_length=3, max_length=500)


class ChunkPatch(BaseModel):
    content: str = Field(min_length=1, max_length=100_000)


class ChunkBatchRequest(BaseModel):
    items: list[dict[str, Any]] = Field(min_length=1, max_length=100)


class RetrieveRequest(BaseModel):
    query: str = Field(min_length=1, max_length=20_000)
    top_k: int | None = Field(default=None, ge=1, le=50)
    score_threshold: float | None = Field(default=None, ge=0.0, le=1.0)


class IndexGenerationCreate(BaseModel):
    reason: str = Field(default="Rebuild requested", min_length=3, max_length=500)


def _safe_name(value: str) -> str:
    return re.sub(r"[^A-Za-z0-9._-]+", "-", value).strip(".-")[:160] or "document"


def _etag(version: int) -> str:
    return f'W/"{version}"'


def _assert_if_match(if_match: str | None, version: int) -> None:
    if not if_match:
        raise HTTPException(status_code=428, detail="If-Match is required")
    accepted = {"*", str(version), _etag(version), f'"{version}"'}
    if if_match.strip() not in accepted:
        raise HTTPException(status_code=412, detail="Resource version does not match If-Match")


def _require_dataset_capability(actor: Actor, action: str) -> None:
    if not capability_allows(actor.capabilities, f"dataset:{action}"):
        raise HTTPException(status_code=403, detail=f"Missing capability: dataset:{action}")


def _normalize_retrieval_config(value: dict[str, Any] | None) -> dict[str, Any]:
    data = {**DEFAULT_RETRIEVAL_CONFIG, **(value or {})}
    try:
        parsed = RetrievalConfigUpsert.model_validate(data)
    except ValueError as exc:
        raise HTTPException(status_code=422, detail=str(exc)) from exc
    return parsed.normalized()


async def _dataset(conn, dataset_id: str, workspace_id: str, *, include_deleted: bool = False) -> dict[str, Any]:
    suffix = "" if include_deleted else " AND status <> 'deleted'"
    result = await conn.execute(
        f"SELECT * FROM pf_dataset WHERE id=%s AND workspace_id=%s{suffix}",
        (dataset_id, workspace_id),
    )
    row = await result.fetchone()
    if not row:
        raise HTTPException(status_code=404, detail="Dataset not found")
    return row


async def _document(conn, dataset_id: str, document_id: str, workspace_id: str, *, include_deleted: bool = False) -> dict[str, Any]:
    suffix = "" if include_deleted else " AND status <> 'deleted'"
    result = await conn.execute(
        f"SELECT * FROM pf_document WHERE id=%s AND dataset_id=%s AND workspace_id=%s{suffix}",
        (document_id, dataset_id, workspace_id),
    )
    row = await result.fetchone()
    if not row:
        raise HTTPException(status_code=404, detail="Document not found")
    return row


async def _outbox(conn, event_type: str, workspace_id: str, payload: dict[str, Any]) -> None:
    await conn.execute(
        """
        INSERT INTO ops_outbox(id,event_type,workspace_id,trace_id,payload)
        VALUES (%s,%s,%s,%s,%s::jsonb)
        """,
        (new_id("out"), event_type, workspace_id, payload.get("operation_id"), json_dumps(payload)),
    )


async def _refresh_dataset_stats(conn, dataset_id: str) -> None:
    await conn.execute(
        """
        UPDATE pf_dataset SET stats=jsonb_build_object(
          'document_count',(SELECT count(*) FROM pf_document WHERE dataset_id=%s AND status <> 'deleted'),
          'chunk_count',(SELECT count(*) FROM pf_chunk WHERE dataset_id=%s)
        ),updated_at=now() WHERE id=%s
        """,
        (dataset_id, dataset_id, dataset_id),
    )


def _mime_for_name(name: str, declared: str | None = None) -> str:
    normalized = (declared or "").split(";", 1)[0].strip().lower()
    suffix = Path(name).suffix.lower()
    if suffix in SUPPORTED_SUFFIXES:
        return SUPPORTED_SUFFIXES[suffix]
    if normalized in TEXT_MIME_TYPES | IMAGE_MIME_TYPES | set(SUPPORTED_SUFFIXES.values()):
        return normalized
    guessed = mimetypes.guess_type(name)[0]
    if guessed in TEXT_MIME_TYPES | IMAGE_MIME_TYPES | set(SUPPORTED_SUFFIXES.values()):
        return guessed
    raise HTTPException(status_code=415, detail="E03001 Document format is not supported")


def _check_content(content: bytes) -> None:
    if len(content) > MAX_DOCUMENT_BYTES:
        raise HTTPException(status_code=413, detail="E03002 Document exceeds 100 MiB")
    if b"EICAR-STANDARD-ANTIVIRUS-TEST-FILE" in content.upper():
        raise HTTPException(status_code=422, detail="Document failed malware scanning")


def _check_zip_expansion(content: bytes) -> None:
    try:
        with zipfile.ZipFile(io.BytesIO(content)) as archive:
            total = sum(item.file_size for item in archive.infolist())
    except zipfile.BadZipFile as exc:
        raise RagProcessingError("E03001 Office document is invalid") from exc
    if total > MAX_DOCUMENT_BYTES * 5:
        raise RagProcessingError("E03001 Office document expands beyond the safe processing limit")


def _extract_text(content: bytes, mime: str, name: str) -> list[tuple[str, dict[str, Any]]]:
    try:
        if mime in TEXT_MIME_TYPES:
            text = content.decode("utf-8-sig", errors="replace")
            if mime == "text/html":
                text = BeautifulSoup(text, "html.parser").get_text("\n")
            return [(text, {"source_name": name})]
        if mime == "application/pdf":
            reader = PdfReader(io.BytesIO(content))
            return [
                (page.extract_text() or "", {"source_name": name, "page": index + 1})
                for index, page in enumerate(reader.pages)
            ]
        if mime == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            _check_zip_expansion(content)
            document = DocxDocument(io.BytesIO(content))
            return [("\n".join(paragraph.text for paragraph in document.paragraphs), {"source_name": name})]
        if mime == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            _check_zip_expansion(content)
            deck = Presentation(io.BytesIO(content))
            pages: list[tuple[str, dict[str, Any]]] = []
            for index, slide in enumerate(deck.slides):
                text = "\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text"))
                pages.append((text, {"source_name": name, "slide": index + 1}))
            return pages
        if mime == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
            _check_zip_expansion(content)
            workbook = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
            pages = []
            for sheet in workbook.worksheets:
                rows = ["\t".join(str(value) if value is not None else "" for value in row) for row in sheet.iter_rows(values_only=True)]
                pages.append(("\n".join(rows), {"source_name": name, "sheet": sheet.title}))
            return pages
        if mime in IMAGE_MIME_TYPES:
            image = Image.open(io.BytesIO(content))
            return [(image_to_string(image, lang="eng+chi_sim"), {"source_name": name, "ocr": True})]
    except RagProcessingError:
        raise
    except Exception as exc:
        raise RagProcessingError(f"E03001 Unable to parse document: {type(exc).__name__}") from exc
    raise RagProcessingError("E03001 Document format is not supported")


def _clean_text(value: str) -> str:
    text = value.replace("\x00", "").replace("\r\n", "\n").replace("\r", "\n")
    text = re.sub(r"[\t \u00a0]+", " ", text)
    lines = [line.strip() for line in text.split("\n")]
    duplicate_counts: dict[str, int] = {}
    for line in lines:
        if len(line) >= 8:
            duplicate_counts[line] = duplicate_counts.get(line, 0) + 1
    cleaned = [line for line in lines if not (len(line) >= 8 and duplicate_counts.get(line, 0) > 3)]
    return "\n".join(cleaned).strip()[:MAX_EXTRACTED_CHARS]


def _title_path(value: str, current: list[str]) -> list[str]:
    match = re.match(r"^(#{1,6})\s+(.+)$", value.strip())
    if not match:
        return current
    level = len(match.group(1))
    return current[: level - 1] + [match.group(2).strip()[:240]]


def _split_text(value: str, metadata: dict[str, Any]) -> list[tuple[str, dict[str, Any]]]:
    target_chars = 2_048
    overlap_chars = 256
    paragraphs = [item.strip() for item in re.split(r"\n{2,}", value) if item.strip()]
    chunks: list[tuple[str, dict[str, Any]]] = []
    current = ""
    titles: list[str] = []
    for paragraph in paragraphs:
        titles = _title_path(paragraph, titles)
        parts = [paragraph[index:index + target_chars] for index in range(0, max(len(paragraph), 1), target_chars)]
        for part in parts:
            candidate = f"{current}\n\n{part}".strip() if current else part
            if current and len(candidate) > target_chars:
                chunk_metadata = dict(metadata)
                if titles:
                    chunk_metadata["title_path"] = titles
                chunks.append((current, chunk_metadata))
                current = f"{current[-overlap_chars:]}\n\n{part}".strip()
            else:
                current = candidate
    if current:
        chunk_metadata = dict(metadata)
        if titles:
            chunk_metadata["title_path"] = titles
        chunks.append((current, chunk_metadata))
    return [(text, item) for text, item in chunks if text.strip()]


def _normalize_vector(values: Iterable[Any]) -> list[float]:
    vector = [float(value) for value in values]
    if not vector or len(vector) > 4000 or any(not math.isfinite(value) for value in vector):
        raise RagProcessingError("Gateway returned an invalid embedding")
    magnitude = math.sqrt(sum(value * value for value in vector))
    if magnitude == 0:
        raise RagProcessingError("Gateway returned a zero embedding")
    return [value / magnitude for value in vector]


def _vector_literal(vector: list[float]) -> str:
    return "[" + ",".join(f"{value:.12g}" for value in vector) + "]"


async def embed_text(workspace_id: str, model: str, text: str, *, request_id: str) -> list[float]:
    headers = {
        "X-Internal-Token": settings.internal_token,
        "X-Workspace-ID": workspace_id,
        "X-Request-ID": request_id,
    }
    try:
        async with httpx.AsyncClient(timeout=90) as client:
            response = await client.post(
                f"{settings.gateway_url.rstrip('/')}/v1/embeddings",
                headers=headers,
                json={"model": model, "input": text},
            )
        if response.status_code >= 300:
            raise RagProcessingError(f"Gateway embedding failed ({response.status_code})")
        payload = response.json()
        embedding = payload["data"][0]["embedding"]
        if isinstance(embedding, str):
            raw = __import__("base64").b64decode(embedding)
            embedding = struct.unpack(f"<{len(raw) // 4}f", raw)
        return _normalize_vector(embedding)
    except RagProcessingError:
        raise
    except Exception as exc:
        raise RagProcessingError(f"Gateway embedding failed: {type(exc).__name__}") from exc


async def _ensure_embedding_index(conn, dimension: int) -> None:
    if not 1 <= dimension <= 4000:
        raise RagProcessingError("Embedding dimension is outside the supported range")
    sql_type = f"vector({dimension})" if dimension <= 2000 else f"halfvec({dimension})"
    opclass = "vector_cosine_ops" if dimension <= 2000 else "halfvec_cosine_ops"
    await conn.execute(
        f"""
        CREATE INDEX IF NOT EXISTS idx_pf_chunk_embedding_{dimension}
        ON pf_chunk USING hnsw ((embedding::{sql_type}) {opclass})
        WHERE embedding_dimension = {dimension}
        """
    )


async def _set_document_status(document_id: str, workspace_id: str, status: str, *, error: str | None = None) -> None:
    async with pool.connection() as conn:
        await conn.execute(
            "UPDATE pf_document SET status=%s,error=%s,updated_at=now(),version=version+1 WHERE id=%s AND workspace_id=%s",
            (status, error, document_id, workspace_id),
        )
        await conn.commit()


async def _assert_not_cancelled(job: ClaimedJob) -> None:
    async with pool.connection() as conn:
        result = await conn.execute("SELECT status FROM ops_async_operation WHERE id=%s", (job.operation_id,))
        operation = await result.fetchone()
    if not operation or operation["status"] in {"cancel_requested", "cancelled"}:
        raise RagJobCancelled("RAG operation was cancelled")


async def _load_document_for_job(document_id: str, workspace_id: str) -> tuple[dict[str, Any], dict[str, Any], dict[str, Any]]:
    async with pool.connection() as conn:
        result = await conn.execute(
            """
            SELECT d.*, ds.embedding_model,ds.embedding_profile,ds.active_generation_id,ds.status AS dataset_status,
                   g.id AS generation_id,g.status AS generation_status
            FROM pf_document d
            JOIN pf_dataset ds ON ds.id=d.dataset_id
            LEFT JOIN pf_index_generation g ON g.id=ds.active_generation_id
            WHERE d.id=%s AND d.workspace_id=%s
            """,
            (document_id, workspace_id),
        )
        row = await result.fetchone()
    if not row or row["status"] in {"deleted", "deleting", "cancelled"} or row["dataset_status"] != "active":
        raise RagJobCancelled("Document is no longer eligible for processing")
    if not row["generation_id"]:
        raise RagProcessingError("Dataset does not have an active index generation")
    return row, {"id": row["dataset_id"], "embedding_model": row["embedding_model"]}, {"id": row["generation_id"]}


async def process_document_job(job: ClaimedJob, *, generation_id: str | None = None) -> dict[str, Any]:
    document, dataset, active_generation = await _load_document_for_job(job.payload["document_id"], job.workspace_id)
    target_generation = generation_id or active_generation["id"]
    await _assert_not_cancelled(job)
    await _set_document_status(document["id"], job.workspace_id, "parsing")
    try:
        content = await get_object(DATASET_BUCKET, document["s3_key"])
        extracted = await asyncio.to_thread(_extract_text, content, document["mime"], document["name"])
        pieces: list[tuple[str, dict[str, Any]]] = []
        for value, metadata in extracted:
            cleaned = _clean_text(value)
            if cleaned:
                pieces.extend(_split_text(cleaned, metadata))
        if not pieces:
            raise RagProcessingError("E03001 No readable text was extracted from the document")
        await _assert_not_cancelled(job)
        await _set_document_status(document["id"], job.workspace_id, "chunking")
        await _set_document_status(document["id"], job.workspace_id, "embedding")
        vectors: list[list[float]] = []
        for index, (text, _) in enumerate(pieces):
            await _assert_not_cancelled(job)
            vectors.append(
                await embed_text(
                    job.workspace_id,
                    dataset["embedding_model"],
                    text,
                    request_id=f"rag_{job.operation_id}_{index}",
                )
            )
        dimensions = {len(vector) for vector in vectors}
        if len(dimensions) != 1:
            raise RagProcessingError("Embedding model returned inconsistent vector dimensions")
        dimension = dimensions.pop()
        async with pool.connection() as conn:
            async with conn.transaction():
                await conn.execute(
                    "DELETE FROM pf_chunk WHERE document_id=%s AND generation_id=%s",
                    (document["id"], target_generation),
                )
                for position, ((text, metadata), vector) in enumerate(zip(pieces, vectors, strict=True)):
                    await conn.execute(
                        """
                        INSERT INTO pf_chunk(id,document_id,dataset_id,generation_id,workspace_id,content,content_sha256,
                          token_count,position,metadata,tsv,embedding,embedding_dimension)
                        VALUES (%s,%s,%s,%s,%s,%s,%s,%s,%s,%s::jsonb,to_tsvector('simple',%s),%s::vector,%s)
                        """,
                        (
                            new_id("chk"), document["id"], document["dataset_id"], target_generation,
                            job.workspace_id, text, hashlib.sha256(text.encode()).hexdigest(),
                            max(1, len(text) // 4), position, json_dumps(metadata), text,
                            _vector_literal(vector), dimension,
                        ),
                    )
                await _ensure_embedding_index(conn, dimension)
                await conn.execute(
                    "UPDATE pf_document SET status='indexed',error=NULL,chunk_count=%s,indexed_at=now(),updated_at=now(),version=version+1 WHERE id=%s",
                    (len(pieces), document["id"]),
                )
                await conn.execute(
                    """
                    UPDATE pf_index_generation SET
                      document_count=(SELECT count(*) FROM pf_document WHERE dataset_id=%s AND status='indexed'),
                      chunk_count=(SELECT count(*) FROM pf_chunk WHERE dataset_id=%s AND generation_id=%s),
                      embedding_profile=jsonb_build_object('provider','gateway','model',%s::text,'dimension',%s::int,'normalization','l2','version',1)
                      WHERE id=%s
                    """,
                    (document["dataset_id"], document["dataset_id"], target_generation, dataset["embedding_model"], dimension, target_generation),
                )
                await conn.execute(
                    "UPDATE pf_dataset SET embedding_profile=jsonb_build_object('provider','gateway','model',%s::text,'dimension',%s::int,'normalization','l2','version',1),updated_at=now() WHERE id=%s AND active_generation_id=%s",
                    (dataset["embedding_model"], dimension, document["dataset_id"], active_generation["id"]),
                )
                await _refresh_dataset_stats(conn, document["dataset_id"])
                await _outbox(
                    conn,
                    "rag.index.activated.v1" if target_generation == active_generation["id"] else "rag.step.requested.v1",
                    job.workspace_id,
                    {"operation_id": job.operation_id, "document_id": document["id"], "dataset_id": document["dataset_id"], "generation_id": target_generation, "chunk_count": len(pieces)},
                )
        return {"document_id": document["id"], "dataset_id": document["dataset_id"], "generation_id": target_generation, "chunk_count": len(pieces), "status": "indexed"}
    except RagJobCancelled:
        await _set_document_status(document["id"], job.workspace_id, "cancelled")
        raise
    except Exception as exc:
        await _set_document_status(document["id"], job.workspace_id, "failed", error=str(exc)[:500])
        raise


async def process_chunk_embedding_job(job: ClaimedJob) -> dict[str, Any]:
    chunk_id = job.payload["chunk_id"]
    async with pool.connection() as conn:
        result = await conn.execute(
            """
            SELECT c.*,d.embedding_model FROM pf_chunk c
            JOIN pf_dataset d ON d.id=c.dataset_id
            WHERE c.id=%s AND c.workspace_id=%s
            """,
            (chunk_id, job.workspace_id),
        )
        chunk = await result.fetchone()
    if not chunk:
        raise RagJobCancelled("Chunk was removed before embedding")
    vector = await embed_text(job.workspace_id, chunk["embedding_model"], chunk["content"], request_id=f"rag_{job.operation_id}")
    async with pool.connection() as conn:
        async with conn.transaction():
            await conn.execute(
                "UPDATE pf_chunk SET embedding=%s::vector,embedding_dimension=%s,tsv=to_tsvector('simple',content),updated_at=now(),version=version+1 WHERE id=%s",
                (_vector_literal(vector), len(vector), chunk_id),
            )
            await _ensure_embedding_index(conn, len(vector))
    return {"chunk_id": chunk_id, "status": "indexed"}


async def process_document_delete_job(job: ClaimedJob) -> dict[str, Any]:
    document_id = job.payload["document_id"]
    async with pool.connection() as conn:
        result = await conn.execute(
            "SELECT id,dataset_id,s3_key,status FROM pf_document WHERE id=%s AND workspace_id=%s",
            (document_id, job.workspace_id),
        )
        document = await result.fetchone()
    if not document or document["status"] == "deleted":
        return {"document_id": document_id, "status": "deleted", "already_deleted": True}
    await _assert_not_cancelled(job)
    async with pool.connection() as conn:
        async with conn.transaction():
            await conn.execute(
                "UPDATE pf_document SET status='deleted',deleted_at=now(),updated_at=now(),version=version+1 WHERE id=%s",
                (document_id,),
            )
            await _refresh_dataset_stats(conn, document["dataset_id"])
            await _outbox(
                conn,
                "rag.step.requested.v1",
                job.workspace_id,
                {"operation_id": job.operation_id, "document_id": document_id, "dataset_id": document["dataset_id"], "action": "deleted"},
            )
    return {"document_id": document_id, "status": "deleted"}


async def process_dataset_delete_job(job: ClaimedJob) -> dict[str, Any]:
    dataset_id = job.payload["dataset_id"]
    async with pool.connection() as conn:
        result = await conn.execute(
            "SELECT id,status FROM pf_dataset WHERE id=%s AND workspace_id=%s",
            (dataset_id, job.workspace_id),
        )
        dataset = await result.fetchone()
        if not dataset or dataset["status"] == "deleted":
            return {"dataset_id": dataset_id, "status": "deleted", "already_deleted": True}
        documents = await conn.execute(
            "SELECT id FROM pf_document WHERE dataset_id=%s AND status <> 'deleted'",
            (dataset_id,),
        )
        rows = await documents.fetchall()
    await _assert_not_cancelled(job)
    async with pool.connection() as conn:
        async with conn.transaction():
            await conn.execute("UPDATE pf_document SET status='deleted',deleted_at=now(),updated_at=now(),version=version+1 WHERE dataset_id=%s AND status <> 'deleted'", (dataset_id,))
            await conn.execute("UPDATE pf_dataset SET status='deleted',deleted_at=now(),updated_at=now(),version=version+1 WHERE id=%s", (dataset_id,))
            await _outbox(
                conn,
                "rag.step.requested.v1",
                job.workspace_id,
                {"operation_id": job.operation_id, "dataset_id": dataset_id, "action": "deleted", "document_count": len(rows)},
            )
    return {"dataset_id": dataset_id, "status": "deleted", "document_count": len(rows)}


async def process_dataset_rebuild_job(job: ClaimedJob) -> dict[str, Any]:
    dataset_id = job.payload["dataset_id"]
    generation_id = job.payload["generation_id"]
    async with pool.connection() as conn:
        result = await conn.execute(
            "SELECT id,status FROM pf_dataset WHERE id=%s AND workspace_id=%s",
            (dataset_id, job.workspace_id),
        )
        dataset = await result.fetchone()
        if not dataset or dataset["status"] != "active":
            raise RagJobCancelled("Dataset is unavailable for rebuilding")
        documents = await conn.execute(
            "SELECT id FROM pf_document WHERE dataset_id=%s AND workspace_id=%s AND status='indexed' ORDER BY created_at",
            (dataset_id, job.workspace_id),
        )
        rows = await documents.fetchall()
    completed = 0
    for item in rows:
        await _assert_not_cancelled(job)
        nested_payload = {**job.payload, "document_id": item["id"]}
        nested_job = ClaimedJob(
            id=job.id,
            operation_id=job.operation_id,
            workspace_id=job.workspace_id,
            job_type="rag.document.process",
            payload=nested_payload,
            attempt_count=job.attempt_count,
            max_attempts=job.max_attempts,
            lease_token=job.lease_token,
        )
        await process_document_job(nested_job, generation_id=generation_id)
        completed += 1
    async with pool.connection() as conn:
        async with conn.transaction():
            await conn.execute(
                "UPDATE pf_index_generation SET status='ready',completed_at=now(),document_count=%s,chunk_count=(SELECT count(*) FROM pf_chunk WHERE generation_id=%s) WHERE id=%s",
                (completed, generation_id, generation_id),
            )
            await _outbox(
                conn,
                "rag.step.requested.v1",
                job.workspace_id,
                {"operation_id": job.operation_id, "dataset_id": dataset_id, "generation_id": generation_id, "action": "rebuild_completed", "document_count": completed},
            )
    return {"dataset_id": dataset_id, "generation_id": generation_id, "document_count": completed, "status": "ready"}


async def process_rag_job(job: ClaimedJob) -> dict[str, Any]:
    if job.job_type.startswith("rag.eval."):
        from workama_platform.modules.rag_eval import process_eval_job

        return await process_eval_job(job)
    if job.job_type.startswith("kb.eval."):
        # T-M3-003: 知识库评测集/标注回流作业，由 knowledge_eval 模块处理
        from workama_platform.modules.knowledge_eval import process_kb_eval_job

        return await process_kb_eval_job(job)
    if job.job_type == "rag.document.process":
        return await process_document_job(job)
    if job.job_type == "rag.document.delete":
        return await process_document_delete_job(job)
    if job.job_type == "rag.dataset.delete":
        return await process_dataset_delete_job(job)
    if job.job_type == "rag.chunk.embed":
        return await process_chunk_embedding_job(job)
    if job.job_type == "rag.dataset.rebuild":
        return await process_dataset_rebuild_job(job)
    raise ValueError(f"Unknown RAG job type: {job.job_type}")


def _dataset_summary(row: dict[str, Any]) -> dict[str, Any]:
    return {
        "id": row["id"],
        "name": row["name"],
        "description": row["description"],
        "embedding_model": row["embedding_model"],
        "embedding_profile": row["embedding_profile"],
        "retrieval_config": row["retrieval_config"],
        "stats": row["stats"],
        "active_generation_id": row["active_generation_id"],
        "status": row["status"],
        "version": row["version"],
        "created_at": row["created_at"],
        "updated_at": row["updated_at"],
        "deleted_at": row["deleted_at"],
    }


def _document_summary(row: dict[str, Any]) -> dict[str, Any]:
    return {
        "id": row["id"],
        "dataset_id": row["dataset_id"],
        "name": row["name"],
        "source": row["source"],
        "source_url": row["source_url"],
        "mime": row["mime"],
        "size_bytes": row["size_bytes"],
        "content_sha256": row["content_sha256"],
        "status": row["status"],
        "error": row["error"],
        "chunk_count": row["chunk_count"],
        "version": row["version"],
        "created_at": row["created_at"],
        "updated_at": row["updated_at"],
        "indexed_at": row["indexed_at"],
        "deleted_at": row["deleted_at"],
    }


def _chunk_summary(row: dict[str, Any]) -> dict[str, Any]:
    return {
        "id": row["id"],
        "document_id": row["document_id"],
        "dataset_id": row["dataset_id"],
        "generation_id": row["generation_id"],
        "content": row["content"],
        "token_count": row["token_count"],
        "position": row["position"],
        "parent_id": row["parent_id"],
        "metadata": row["metadata"],
        "version": row["version"],
        "created_at": row["created_at"],
        "updated_at": row["updated_at"],
    }


async def _submit_rag_operation(
    conn, *, actor: Actor, operation_type: str, job_type: str, payload: dict[str, Any],
    idempotency_key: str, max_attempts: int = 3,
) -> dict[str, Any]:
    try:
        operation = await submit_operation(
            conn,
            operation_type=operation_type,
            workspace_id=actor.workspace_id,
            org_id=actor.org_id,
            actor_id=actor.user_id,
            actor_role=actor.role,
            idempotency_key=idempotency_key,
            payload=payload,
            job_type=job_type,
            queue="rag",
            max_attempts=max_attempts,
            priority=100,
            cancellable=True,
        )
    except IdempotencyConflict as exc:
        raise HTTPException(status_code=409, detail="E00008 Idempotency key was already used with different input") from exc
    await _outbox(conn, "rag.document.accepted.v1" if job_type == "rag.document.process" else "rag.step.requested.v1", actor.workspace_id, {"operation_id": operation["id"], **payload, "job_type": job_type})
    return operation


@router.get("/datasets")
async def list_datasets(actor: Annotated[Actor, Depends(get_actor)], limit: int = 50):
    _require_dataset_capability(actor, "read")
    async with pool.connection() as conn:
        result = await conn.execute(
            "SELECT * FROM pf_dataset WHERE workspace_id=%s AND status <> 'deleted' ORDER BY updated_at DESC LIMIT %s",
            (actor.workspace_id, max(1, min(limit, 200))),
        )
        data = [_dataset_summary(row) for row in await result.fetchall()]
    # Contract《720》listDatasets: ListQuery -> ListResponse<DatasetDTO>
    return {
        "items": data,
        "data": data,
        "next_cursor": None,
        "has_more": False,
        "meta": {"request_id": None},
    }


@router.post("/datasets", status_code=201)
async def create_dataset(body: DatasetCreate, response: Response, actor: Annotated[Actor, Depends(get_actor)]):
    _require_dataset_capability(actor, "create")
    dataset_id = new_id("dts")
    generation_id = new_id("idx")
    config = _normalize_retrieval_config(body.retrieval_config)
    async with pool.connection() as conn:
        async with conn.transaction():
            try:
                result = await conn.execute(
                    """
                    INSERT INTO pf_dataset(id,org_id,workspace_id,name,description,embedding_model,retrieval_config,active_generation_id,created_by)
                    VALUES (%s,%s,%s,%s,%s,%s,%s::jsonb,%s,%s) RETURNING *
                    """,
                    (dataset_id, actor.org_id, actor.workspace_id, body.name, body.description, body.embedding_model, json_dumps(config), generation_id, actor.user_id),
                )
                row = await result.fetchone()
            except Exception as exc:
                if "unique" in str(exc).lower():
                    raise HTTPException(status_code=409, detail="A dataset with this name already exists") from exc
                raise
            await conn.execute(
                """
                INSERT INTO pf_index_generation(id,dataset_id,workspace_id,generation,embedding_profile,status,created_by,activated_at,completed_at)
                VALUES (%s,%s,%s,1,%s::jsonb,'active',%s,now(),now())
                """,
                (generation_id, dataset_id, actor.workspace_id, json_dumps({"provider": "gateway", "model": body.embedding_model, "normalization": "l2", "version": 1}), actor.user_id),
            )
    response.headers["ETag"] = _etag(row["version"])
    return _dataset_summary(row)


@router.get("/datasets/{dataset_id}")
async def get_dataset(dataset_id: str, response: Response, actor: Annotated[Actor, Depends(get_actor)]):
    _require_dataset_capability(actor, "read")
    async with pool.connection() as conn:
        row = await _dataset(conn, dataset_id, actor.workspace_id)
    response.headers["ETag"] = _etag(row["version"])
    return _dataset_summary(row)


@router.patch("/datasets/{dataset_id}")
async def update_dataset(
    dataset_id: str,
    body: DatasetPatch,
    response: Response,
    actor: Annotated[Actor, Depends(get_actor)],
    if_match: Annotated[str | None, Header(alias="If-Match")] = None,
):
    _require_dataset_capability(actor, "write")
    if body.name is None and body.description is None:
        raise HTTPException(status_code=422, detail="At least one mutable field is required")
    async with pool.connection() as conn:
        async with conn.transaction():
            current = await _dataset(conn, dataset_id, actor.workspace_id)
            _assert_if_match(if_match, current["version"])
            result = await conn.execute(
                """
                UPDATE pf_dataset SET name=COALESCE(%s,name),description=COALESCE(%s,description),version=version+1,updated_at=now()
                WHERE id=%s AND workspace_id=%s AND version=%s RETURNING *
                """,
                (body.name, body.description, dataset_id, actor.workspace_id, current["version"]),
            )
            row = await result.fetchone()
    if not row:
        raise HTTPException(status_code=412, detail="Resource version changed")
    response.headers["ETag"] = _etag(row["version"])
    return _dataset_summary(row)


@router.delete("/datasets/{dataset_id}", status_code=202)
async def delete_dataset(
    dataset_id: str,
    body: DeleteReason,
    actor: Annotated[Actor, Depends(get_actor)],
    if_match: Annotated[str | None, Header(alias="If-Match")] = None,
    idempotency_key: Annotated[str | None, Header(alias="Idempotency-Key")] = None,
):
    _require_dataset_capability(actor, "delete")
    async with pool.connection() as conn:
        async with conn.transaction():
            dataset = await _dataset(conn, dataset_id, actor.workspace_id)
            _assert_if_match(if_match, dataset["version"])
            await conn.execute("UPDATE pf_dataset SET status='deleting',updated_at=now() WHERE id=%s", (dataset_id,))
            operation = await _submit_rag_operation(
                conn,
                actor=actor,
                operation_type="rag.dataset.delete",
                job_type="rag.dataset.delete",
                payload={"dataset_id": dataset_id, "reason": body.reason},
                idempotency_key=idempotency_key or f"dataset-delete:{dataset_id}:{dataset['version']}",
            )
    return {"operation": operation, "dataset_id": dataset_id, "status": "deleting"}


@router.post("/datasets/{dataset_id}/restore")
async def restore_dataset(dataset_id: str, body: RestoreRequest, actor: Annotated[Actor, Depends(get_actor)]):
    _require_dataset_capability(actor, "restore")
    async with pool.connection() as conn:
        result = await conn.execute(
            "UPDATE pf_dataset SET status='active',deleted_at=NULL,updated_at=now(),version=version+1 WHERE id=%s AND workspace_id=%s AND status='deleted' RETURNING *",
            (dataset_id, actor.workspace_id),
        )
        row = await result.fetchone()
        if not row:
            raise HTTPException(status_code=404, detail="Deleted dataset not found")
        await conn.commit()
    return _dataset_summary(row)


async def _read_document_upload(request: Request) -> tuple[str, str, str | None, bytes, str]:
    content_type = request.headers.get("content-type", "")
    source_url: str | None = None
    declared_mime = ""
    if content_type.startswith("multipart/"):
        form = await request.form()
        uploaded = form.get("file")
        source_url = str(form.get("source_url") or "").strip() or None
        preferred_name = str(form.get("name") or "").strip()
        if uploaded is not None and hasattr(uploaded, "read"):
            content = await uploaded.read(MAX_DOCUMENT_BYTES + 1)
            name = preferred_name or getattr(uploaded, "filename", None) or "upload"
            declared_mime = getattr(uploaded, "content_type", "") or ""
            return name, _mime_for_name(name, declared_mime), None, content, "upload"
        if not source_url:
            raise HTTPException(status_code=422, detail="Provide a file or source_url")
        name = preferred_name
    else:
        try:
            body = await request.json()
        except Exception as exc:
            raise HTTPException(status_code=422, detail="Upload requests must be multipart or contain a source_url JSON object") from exc
        source_url = str(body.get("source_url") or "").strip() or None
        name = str(body.get("name") or "").strip()
        if not source_url:
            raise HTTPException(status_code=422, detail="source_url is required for JSON document creation")
    validation = await validate_resolved_outbound_url(source_url)
    if not validation.allowed:
        raise HTTPException(status_code=422, detail=f"E03001 Unsafe source URL: {validation.reason}")
    try:
        async with httpx.AsyncClient(timeout=30, follow_redirects=False) as client:
            response = await client.get(source_url, headers={"User-Agent": "WorkAMA-RAG/1.0"})
        if response.is_redirect:
            raise HTTPException(status_code=422, detail="E03001 Source URL redirects are not allowed")
        response.raise_for_status()
    except HTTPException:
        raise
    except Exception as exc:
        raise HTTPException(status_code=422, detail=f"E03001 Unable to fetch source URL: {type(exc).__name__}") from exc
    path_name = Path(urlsplit(source_url).path).name
    name = name or path_name or "remote-document"
    declared_mime = response.headers.get("content-type", "")
    return name, _mime_for_name(name, declared_mime), source_url, response.content, "url"


@router.get("/datasets/{dataset_id}/documents")
async def list_dataset_documents(
    dataset_id: str,
    actor: Annotated[Actor, Depends(get_actor)],
    limit: int = 100,
    include_deleted: bool = False,
):
    _require_dataset_capability(actor, "read")
    async with pool.connection() as conn:
        await _dataset(conn, dataset_id, actor.workspace_id)
        status_clause = "" if include_deleted else " AND status <> 'deleted'"
        result = await conn.execute(
            f"SELECT * FROM pf_document WHERE dataset_id=%s AND workspace_id=%s{status_clause} ORDER BY created_at DESC LIMIT %s",
            (dataset_id, actor.workspace_id, max(1, min(limit, 200))),
        )
        data = [_document_summary(row) for row in await result.fetchall()]
    # Contract《720》listDatasetDocuments: ListQuery -> ListResponse<DocumentDTO>
    return {
        "items": data,
        "data": data,
        "next_cursor": None,
        "has_more": False,
        "meta": {"request_id": None},
    }


@router.post("/datasets/{dataset_id}/documents", status_code=202)
async def create_dataset_document(
    dataset_id: str,
    request: Request,
    actor: Annotated[Actor, Depends(get_actor)],
    idempotency_key: Annotated[str | None, Header(alias="Idempotency-Key")] = None,
):
    _require_dataset_capability(actor, "create")
    name, mime, source_url, content, source = await _read_document_upload(request)
    _check_content(content)
    digest = hashlib.sha256(content).hexdigest()
    document_id = new_id("doc")
    key = f"datasets/{actor.workspace_id}/{dataset_id}/{document_id}/{_safe_name(name)}"
    await put_object(DATASET_BUCKET, key, content)
    try:
        async with pool.connection() as conn:
            async with conn.transaction():
                await _dataset(conn, dataset_id, actor.workspace_id)
                existing_result = await conn.execute(
                    "SELECT * FROM pf_document WHERE dataset_id=%s AND content_sha256=%s AND status <> 'deleted' ORDER BY created_at DESC LIMIT 1",
                    (dataset_id, digest),
                )
                existing = await existing_result.fetchone()
                if existing:
                    # The object was uploaded before the database dedupe check;
                    # it is not referenced by the existing document and must be
                    # removed to avoid orphaning tenant data in object storage.
                    await delete_object(DATASET_BUCKET, key)
                    operation = await _submit_rag_operation(
                        conn,
                        actor=actor,
                        operation_type="rag.document.process",
                        job_type="rag.document.process",
                        payload={"document_id": existing["id"], "dataset_id": dataset_id},
                        idempotency_key=idempotency_key or f"document-process:{existing['id']}:{existing['version']}",
                    )
                    return {"document": _document_summary(existing), "operation": operation, "deduplicated": True}
                result = await conn.execute(
                    """
                    INSERT INTO pf_document(id,dataset_id,workspace_id,name,source,source_url,s3_key,mime,size_bytes,content_sha256,created_by)
                    VALUES (%s,%s,%s,%s,%s,%s,%s,%s,%s,%s,%s) RETURNING *
                    """,
                    (document_id, dataset_id, actor.workspace_id, name[:240], source, source_url, key, mime, len(content), digest, actor.user_id),
                )
                document = await result.fetchone()
                operation = await _submit_rag_operation(
                    conn,
                    actor=actor,
                    operation_type="rag.document.process",
                    job_type="rag.document.process",
                    payload={"document_id": document_id, "dataset_id": dataset_id},
                    idempotency_key=idempotency_key or f"document-process:{document_id}",
                )
                await _refresh_dataset_stats(conn, dataset_id)
    except Exception:
        await delete_object(DATASET_BUCKET, key)
        raise
    return {"document": _document_summary(document), "operation": operation, "deduplicated": False}


@router.get("/datasets/{dataset_id}/documents/{document_id}")
async def get_dataset_document(dataset_id: str, document_id: str, response: Response, actor: Annotated[Actor, Depends(get_actor)]):
    _require_dataset_capability(actor, "read")
    async with pool.connection() as conn:
        row = await _document(conn, dataset_id, document_id, actor.workspace_id)
    response.headers["ETag"] = _etag(row["version"])
    return _document_summary(row)


@router.delete("/datasets/{dataset_id}/documents/{document_id}", status_code=202)
async def delete_dataset_document(
    dataset_id: str,
    document_id: str,
    body: DeleteReason,
    actor: Annotated[Actor, Depends(get_actor)],
    if_match: Annotated[str | None, Header(alias="If-Match")] = None,
    idempotency_key: Annotated[str | None, Header(alias="Idempotency-Key")] = None,
):
    _require_dataset_capability(actor, "delete")
    async with pool.connection() as conn:
        async with conn.transaction():
            document = await _document(conn, dataset_id, document_id, actor.workspace_id)
            _assert_if_match(if_match, document["version"])
            await conn.execute("UPDATE pf_document SET status='deleting',updated_at=now(),version=version+1 WHERE id=%s", (document_id,))
            operation = await _submit_rag_operation(
                conn,
                actor=actor,
                operation_type="rag.document.delete",
                job_type="rag.document.delete",
                payload={"document_id": document_id, "dataset_id": dataset_id, "reason": body.reason},
                idempotency_key=idempotency_key or f"document-delete:{document_id}:{document['version']}",
            )
    return {"document_id": document_id, "operation": operation, "status": "deleting"}


@router.post("/datasets/{dataset_id}/documents/{document_id}/restore")
async def restore_dataset_document(dataset_id: str, document_id: str, body: RestoreRequest, actor: Annotated[Actor, Depends(get_actor)]):
    _require_dataset_capability(actor, "restore")
    async with pool.connection() as conn:
        result = await conn.execute(
            "UPDATE pf_document SET status='pending',deleted_at=NULL,updated_at=now(),version=version+1 WHERE id=%s AND dataset_id=%s AND workspace_id=%s AND status='deleted' RETURNING *",
            (document_id, dataset_id, actor.workspace_id),
        )
        document = await result.fetchone()
        if not document:
            raise HTTPException(status_code=404, detail="Deleted document not found")
        operation = await _submit_rag_operation(
            conn,
            actor=actor,
            operation_type="rag.document.process",
            job_type="rag.document.process",
            payload={"document_id": document_id, "dataset_id": dataset_id},
            idempotency_key=f"document-restore:{document_id}:{document['version']}",
        )
        await conn.commit()
    return {"document": _document_summary(document), "operation": operation}


@router.post("/datasets/{dataset_id}/documents/{document_id}/retries", status_code=202)
async def retry_dataset_document(
    dataset_id: str,
    document_id: str,
    actor: Annotated[Actor, Depends(get_actor)],
    idempotency_key: Annotated[str | None, Header(alias="Idempotency-Key")] = None,
):
    _require_dataset_capability(actor, "write")
    async with pool.connection() as conn:
        async with conn.transaction():
            document = await _document(conn, dataset_id, document_id, actor.workspace_id)
            if document["status"] not in {"failed", "cancelled", "indexed"}:
                raise HTTPException(status_code=409, detail="E03003 Document is already processing")
            result = await conn.execute(
                "UPDATE pf_document SET status='pending',error=NULL,updated_at=now(),version=version+1 WHERE id=%s RETURNING *",
                (document_id,),
            )
            row = await result.fetchone()
            operation = await _submit_rag_operation(
                conn,
                actor=actor,
                operation_type="rag.document.process",
                job_type="rag.document.process",
                payload={"document_id": document_id, "dataset_id": dataset_id},
                idempotency_key=idempotency_key or f"document-retry:{document_id}:{row['version']}",
            )
    return {"document": _document_summary(row), "operation": operation}


@router.post("/datasets/{dataset_id}/documents/{document_id}/cancel")
async def cancel_dataset_document(dataset_id: str, document_id: str, actor: Annotated[Actor, Depends(get_actor)]):
    _require_dataset_capability(actor, "write")
    async with pool.connection() as conn:
        async with conn.transaction():
            document = await _document(conn, dataset_id, document_id, actor.workspace_id)
            if document["status"] not in {"pending", "parsing", "chunking", "embedding"}:
                raise HTTPException(status_code=409, detail="Document is not processing")
            await conn.execute("UPDATE pf_document SET status='cancelled',updated_at=now(),version=version+1 WHERE id=%s", (document_id,))
            result = await conn.execute(
                "SELECT operation_id FROM ops_job WHERE queue='rag' AND workspace_id=%s AND payload->>'document_id'=%s AND status IN ('queued','running','retry_wait') ORDER BY created_at DESC LIMIT 1",
                (actor.workspace_id, document_id),
            )
            job = await result.fetchone()
            if job:
                await request_cancellation(conn, operation_id=job["operation_id"], workspace_id=actor.workspace_id, reason="Document processing cancelled")
    return {"id": document_id, "status": "cancelled"}


@router.get("/datasets/{dataset_id}/chunks")
async def list_dataset_chunks(
    dataset_id: str,
    actor: Annotated[Actor, Depends(get_actor)],
    document_id: str | None = None,
    limit: int = 100,
):
    _require_dataset_capability(actor, "read")
    async with pool.connection() as conn:
        dataset = await _dataset(conn, dataset_id, actor.workspace_id)
        query = "SELECT * FROM pf_chunk WHERE dataset_id=%s AND workspace_id=%s AND generation_id=%s"
        params: list[Any] = [dataset_id, actor.workspace_id, dataset["active_generation_id"]]
        if document_id:
            query += " AND document_id=%s"
            params.append(document_id)
        query += " ORDER BY document_id,position LIMIT %s"
        params.append(max(1, min(limit, 500)))
        result = await conn.execute(query, params)
        data = [_chunk_summary(row) for row in await result.fetchall()]
    # Contract《720》listDatasetChunks: ListQuery -> ListResponse<ChunkDTO>
    return {
        "items": data,
        "data": data,
        "next_cursor": None,
        "has_more": False,
        "meta": {"request_id": None},
    }


@router.get("/datasets/{dataset_id}/chunks/{chunk_id}")
async def get_dataset_chunk(dataset_id: str, chunk_id: str, response: Response, actor: Annotated[Actor, Depends(get_actor)]):
    _require_dataset_capability(actor, "read")
    async with pool.connection() as conn:
        result = await conn.execute(
            "SELECT * FROM pf_chunk WHERE id=%s AND dataset_id=%s AND workspace_id=%s",
            (chunk_id, dataset_id, actor.workspace_id),
        )
        chunk = await result.fetchone()
    if not chunk:
        raise HTTPException(status_code=404, detail="Chunk not found")
    response.headers["ETag"] = _etag(chunk["version"])
    return _chunk_summary(chunk)


async def _queue_chunk_embedding(
    conn, *, chunk: dict[str, Any], actor: Actor, idempotency_key: str,
) -> dict[str, Any]:
    return await _submit_rag_operation(
        conn,
        actor=actor,
        operation_type="rag.chunk.embed",
        job_type="rag.chunk.embed",
        payload={"chunk_id": chunk["id"], "dataset_id": chunk["dataset_id"], "document_id": chunk["document_id"]},
        idempotency_key=idempotency_key,
    )


@router.patch("/datasets/{dataset_id}/chunks/{chunk_id}", status_code=202)
async def update_dataset_chunk(
    dataset_id: str,
    chunk_id: str,
    body: ChunkPatch,
    actor: Annotated[Actor, Depends(get_actor)],
    if_match: Annotated[str | None, Header(alias="If-Match")] = None,
    idempotency_key: Annotated[str | None, Header(alias="Idempotency-Key")] = None,
):
    _require_dataset_capability(actor, "write")
    async with pool.connection() as conn:
        async with conn.transaction():
            result = await conn.execute(
                "SELECT * FROM pf_chunk WHERE id=%s AND dataset_id=%s AND workspace_id=%s FOR UPDATE",
                (chunk_id, dataset_id, actor.workspace_id),
            )
            current = await result.fetchone()
            if not current:
                raise HTTPException(status_code=404, detail="Chunk not found")
            _assert_if_match(if_match, current["version"])
            updated = await conn.execute(
                """
                UPDATE pf_chunk SET content=%s,content_sha256=%s,token_count=%s,tsv=to_tsvector('simple',%s),
                  version=version+1,updated_at=now() WHERE id=%s AND version=%s RETURNING *
                """,
                (body.content, hashlib.sha256(body.content.encode()).hexdigest(), max(1, len(body.content) // 4), body.content, chunk_id, current["version"]),
            )
            chunk = await updated.fetchone()
            if not chunk:
                raise HTTPException(status_code=412, detail="Chunk changed before update completed")
            operation = await _queue_chunk_embedding(
                conn,
                chunk=chunk,
                actor=actor,
                idempotency_key=idempotency_key or f"chunk-embed:{chunk_id}:{chunk['version']}",
            )
    return {"chunk": _chunk_summary(chunk), "operation": operation}


@router.post("/datasets/{dataset_id}/chunks/batch", status_code=202)
async def batch_update_dataset_chunks(
    dataset_id: str,
    body: ChunkBatchRequest,
    actor: Annotated[Actor, Depends(get_actor)],
    idempotency_key: Annotated[str | None, Header(alias="Idempotency-Key")] = None,
):
    _require_dataset_capability(actor, "write")
    results: list[dict[str, Any]] = []
    async with pool.connection() as conn:
        async with conn.transaction():
            await _dataset(conn, dataset_id, actor.workspace_id)
            for index, item in enumerate(body.items):
                chunk_id = str(item.get("id") or "")
                content = str(item.get("content") or "").strip()
                version = item.get("version")
                if not chunk_id or not content or not isinstance(version, int):
                    results.append({"index": index, "status": "invalid", "error": "id, content and integer version are required"})
                    continue
                lookup = await conn.execute(
                    "SELECT * FROM pf_chunk WHERE id=%s AND dataset_id=%s AND workspace_id=%s FOR UPDATE",
                    (chunk_id, dataset_id, actor.workspace_id),
                )
                current = await lookup.fetchone()
                if not current:
                    results.append({"id": chunk_id, "status": "not_found"})
                    continue
                if current["version"] != version:
                    results.append({"id": chunk_id, "status": "conflict", "version": current["version"]})
                    continue
                updated = await conn.execute(
                    "UPDATE pf_chunk SET content=%s,content_sha256=%s,token_count=%s,tsv=to_tsvector('simple',%s),version=version+1,updated_at=now() WHERE id=%s RETURNING *",
                    (content, hashlib.sha256(content.encode()).hexdigest(), max(1, len(content) // 4), content, chunk_id),
                )
                chunk = await updated.fetchone()
                operation = await _queue_chunk_embedding(
                    conn,
                    chunk=chunk,
                    actor=actor,
                    idempotency_key=f"{idempotency_key or 'chunk-batch'}:{chunk_id}:{chunk['version']}",
                )
                results.append({"id": chunk_id, "status": "accepted", "chunk": _chunk_summary(chunk), "operation": operation})
    return {"items": results}


async def _retrieve_rows(dataset: dict[str, Any], workspace_id: str, query: str, candidate_k: int) -> tuple[list[dict[str, Any]], list[dict[str, Any]]]:
    profile = dataset.get("embedding_profile") or {}
    dimension = profile.get("dimension")
    if not isinstance(dimension, int) or dimension < 1:
        raise HTTPException(status_code=409, detail="E03003 Dataset index is not ready")
    vector = await embed_text(workspace_id, dataset["embedding_model"], query, request_id=f"rag_query_{new_id('q')}")
    if len(vector) != dimension:
        raise HTTPException(status_code=503, detail="Embedding profile dimension does not match the active gateway model")
    literal = _vector_literal(vector)
    sql_type = f"vector({dimension})" if dimension <= 2000 else f"halfvec({dimension})"
    async with pool.connection() as conn:
        keyword_result = await conn.execute(
            """
            SELECT c.*,d.name AS document_name,ts_rank_cd(c.tsv,websearch_to_tsquery('simple',%s)) AS keyword_score
            FROM pf_chunk c JOIN pf_document d ON d.id=c.document_id
            WHERE c.dataset_id=%s AND c.workspace_id=%s AND c.generation_id=%s
              AND d.status='indexed' AND c.tsv @@ websearch_to_tsquery('simple',%s)
            ORDER BY keyword_score DESC,c.position LIMIT %s
            """,
            (query, dataset["id"], workspace_id, dataset["active_generation_id"], query, candidate_k),
        )
        keyword = await keyword_result.fetchall()
        vector_result = await conn.execute(
            f"""
            SELECT c.*,d.name AS document_name,1 - ((c.embedding::{sql_type}) <=> %s::{sql_type}) AS vector_score
            FROM pf_chunk c JOIN pf_document d ON d.id=c.document_id
            WHERE c.dataset_id=%s AND c.workspace_id=%s AND c.generation_id=%s AND c.embedding_dimension=%s
              AND d.status='indexed'
            ORDER BY (c.embedding::{sql_type}) <=> %s::{sql_type} LIMIT %s
            """,
            (literal, dataset["id"], workspace_id, dataset["active_generation_id"], dimension, literal, candidate_k),
        )
        vector_rows = await vector_result.fetchall()
    return keyword, vector_rows


@router.post("/datasets/{dataset_id}/retrieve")
async def retrieve_dataset(dataset_id: str, body: RetrieveRequest, actor: Annotated[Actor, Depends(get_actor)]):
    _require_dataset_capability(actor, "retrieve")
    async with pool.connection() as conn:
        dataset = await _dataset(conn, dataset_id, actor.workspace_id)
        if dataset["status"] != "active" or not dataset["active_generation_id"]:
            raise HTTPException(status_code=409, detail="E03003 Dataset is not retrievable")
        count_result = await conn.execute(
            "SELECT count(*) AS count FROM pf_chunk WHERE dataset_id=%s AND generation_id=%s",
            (dataset_id, dataset["active_generation_id"]),
        )
        if (await count_result.fetchone())["count"] == 0:
            raise HTTPException(status_code=409, detail="E03003 Dataset index is not ready")
    config = _normalize_retrieval_config(dataset["retrieval_config"])
    top_k = body.top_k or config["top_k"]
    threshold = body.score_threshold if body.score_threshold is not None else config["score_threshold"]
    keyword, vector_rows = await _retrieve_rows(dataset, actor.workspace_id, body.query, config["candidate_k"])
    scores: dict[str, dict[str, Any]] = {}
    for rank, row in enumerate(keyword, start=1):
        item = scores.setdefault(row["id"], {"row": dict(row), "keyword_rank": None, "vector_rank": None, "score": 0.0})
        item["keyword_rank"] = rank
        item["score"] += 1 / (config["rrf_k"] + rank)
    for rank, row in enumerate(vector_rows, start=1):
        item = scores.get(row["id"])
        if item is None:
            item = {"row": dict(row), "keyword_rank": None, "vector_rank": None, "score": 0.0}
            scores[row["id"]] = item
        else:
            # The keyword pass may have inserted this chunk first. Preserve its
            # metadata while carrying over the vector score used by threshold
            # filtering and the response DTO.
            item["row"]["vector_score"] = row.get("vector_score")
        item["vector_rank"] = rank
        item["score"] += 1 / (config["rrf_k"] + rank)
    ordered = sorted(scores.values(), key=lambda item: item["score"], reverse=True)
    items = []
    for item in ordered:
        row = item["row"]
        vector_score = float(row.get("vector_score") or 0.0)
        if threshold and vector_score < threshold:
            continue
        items.append(
            {
                "id": row["id"],
                "document_id": row["document_id"],
                "document_name": row["document_name"],
                "content": row["content"],
                "metadata": row["metadata"],
                "rrf_score": item["score"],
                "keyword_rank": item["keyword_rank"],
                "vector_rank": item["vector_rank"],
                "keyword_score": float(row.get("keyword_score") or 0.0),
                "vector_score": vector_score,
            }
        )
        if len(items) >= top_k:
            break
    return {"dataset_id": dataset_id, "query": body.query, "items": items, "config": config}


@router.get("/datasets/{dataset_id}/retrieval-config")
async def get_dataset_retrieval_config(dataset_id: str, response: Response, actor: Annotated[Actor, Depends(get_actor)]):
    _require_dataset_capability(actor, "read")
    async with pool.connection() as conn:
        dataset = await _dataset(conn, dataset_id, actor.workspace_id)
    response.headers["ETag"] = _etag(dataset["version"])
    return {"dataset_id": dataset_id, "config": _normalize_retrieval_config(dataset["retrieval_config"]), "version": dataset["version"]}


@router.put("/datasets/{dataset_id}/retrieval-config")
async def update_dataset_retrieval_config(
    dataset_id: str,
    body: RetrievalConfigUpsert,
    response: Response,
    actor: Annotated[Actor, Depends(get_actor)],
    if_match: Annotated[str | None, Header(alias="If-Match")] = None,
):
    _require_dataset_capability(actor, "write")
    config = body.normalized()
    async with pool.connection() as conn:
        async with conn.transaction():
            dataset = await _dataset(conn, dataset_id, actor.workspace_id)
            _assert_if_match(if_match, dataset["version"])
            result = await conn.execute(
                "UPDATE pf_dataset SET retrieval_config=%s::jsonb,version=version+1,updated_at=now() WHERE id=%s AND version=%s RETURNING *",
                (json_dumps(config), dataset_id, dataset["version"]),
            )
            row = await result.fetchone()
    if not row:
        raise HTTPException(status_code=412, detail="Dataset changed before update completed")
    response.headers["ETag"] = _etag(row["version"])
    return {"dataset_id": dataset_id, "config": row["retrieval_config"], "version": row["version"]}


@router.get("/datasets/{dataset_id}/index-generations")
async def list_index_generations(dataset_id: str, actor: Annotated[Actor, Depends(get_actor)]):
    _require_dataset_capability(actor, "read")
    async with pool.connection() as conn:
        await _dataset(conn, dataset_id, actor.workspace_id)
        result = await conn.execute(
            "SELECT * FROM pf_index_generation WHERE dataset_id=%s AND workspace_id=%s ORDER BY generation DESC",
            (dataset_id, actor.workspace_id),
        )
        data = await result.fetchall()
    # Contract《720》listIndexGenerations: ListQuery -> ListResponse<IndexGenerationDTO>
    return {
        "items": data,
        "data": data,
        "next_cursor": None,
        "has_more": False,
        "meta": {"request_id": None},
    }


@router.get("/datasets/{dataset_id}/index-generations/{generation_id}")
async def get_index_generation(dataset_id: str, generation_id: str, actor: Annotated[Actor, Depends(get_actor)]):
    _require_dataset_capability(actor, "read")
    async with pool.connection() as conn:
        result = await conn.execute(
            "SELECT * FROM pf_index_generation WHERE id=%s AND dataset_id=%s AND workspace_id=%s",
            (generation_id, dataset_id, actor.workspace_id),
        )
        row = await result.fetchone()
    if not row:
        raise HTTPException(status_code=404, detail="Index generation not found")
    return row


@router.post("/datasets/{dataset_id}/index-generations", status_code=202)
async def rebuild_dataset_index(
    dataset_id: str,
    body: IndexGenerationCreate,
    actor: Annotated[Actor, Depends(get_actor)],
    idempotency_key: Annotated[str | None, Header(alias="Idempotency-Key")] = None,
):
    _require_dataset_capability(actor, "write")
    generation_id = new_id("idx")
    async with pool.connection() as conn:
        async with conn.transaction():
            dataset = await _dataset(conn, dataset_id, actor.workspace_id)
            number_result = await conn.execute("SELECT COALESCE(max(generation),0)+1 AS generation FROM pf_index_generation WHERE dataset_id=%s", (dataset_id,))
            generation = (await number_result.fetchone())["generation"]
            await conn.execute(
                "INSERT INTO pf_index_generation(id,dataset_id,workspace_id,generation,embedding_profile,status,created_by) VALUES (%s,%s,%s,%s,%s::jsonb,'building',%s)",
                (generation_id, dataset_id, actor.workspace_id, generation, json_dumps(dataset["embedding_profile"]), actor.user_id),
            )
            operation = await _submit_rag_operation(
                conn,
                actor=actor,
                operation_type="rag.dataset.rebuild",
                job_type="rag.dataset.rebuild",
                payload={"dataset_id": dataset_id, "generation_id": generation_id, "reason": body.reason},
                idempotency_key=idempotency_key or f"dataset-rebuild:{dataset_id}:{generation}",
                max_attempts=2,
            )
    return {"generation_id": generation_id, "operation": operation, "status": "building"}


@router.post("/datasets/{dataset_id}/index-generations/{generation_id}/activate")
async def activate_index_generation(
    dataset_id: str,
    generation_id: str,
    body: RestoreRequest,
    actor: Annotated[Actor, Depends(get_actor)],
    if_match: Annotated[str | None, Header(alias="If-Match")] = None,
):
    _require_dataset_capability(actor, "write")
    async with pool.connection() as conn:
        async with conn.transaction():
            dataset = await _dataset(conn, dataset_id, actor.workspace_id)
            _assert_if_match(if_match, dataset["version"])
            result = await conn.execute(
                "SELECT * FROM pf_index_generation WHERE id=%s AND dataset_id=%s AND workspace_id=%s FOR UPDATE",
                (generation_id, dataset_id, actor.workspace_id),
            )
            generation = await result.fetchone()
            if not generation:
                raise HTTPException(status_code=404, detail="Index generation not found")
            if generation["status"] != "ready":
                raise HTTPException(status_code=409, detail="E03003 Index generation is not ready")
            await conn.execute("UPDATE pf_index_generation SET status='retired' WHERE dataset_id=%s AND status='active'", (dataset_id,))
            await conn.execute("UPDATE pf_index_generation SET status='active',activated_at=now() WHERE id=%s", (generation_id,))
            activated = await conn.execute(
                """
                UPDATE pf_dataset
                SET active_generation_id=%s,embedding_profile=%s::jsonb,version=version+1,updated_at=now()
                WHERE id=%s AND workspace_id=%s AND version=%s
                RETURNING id
                """,
                (
                    generation_id,
                    json_dumps(generation["embedding_profile"] or {}),
                    dataset_id,
                    actor.workspace_id,
                    dataset["version"],
                ),
            )
            if not await activated.fetchone():
                raise HTTPException(status_code=412, detail="Dataset changed before generation activation completed")
            await _outbox(
                conn,
                "rag.index.activated.v1",
                actor.workspace_id,
                {"dataset_id": dataset_id, "generation_id": generation_id, "actor_id": actor.user_id, "reason": body.reason},
            )
    return {"dataset_id": dataset_id, "generation_id": generation_id, "status": "active"}
