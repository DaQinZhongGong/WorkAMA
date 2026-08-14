from __future__ import annotations

import asyncio
import hashlib
import ipaddress
import json
import re
import socket
from dataclasses import dataclass
from datetime import UTC, datetime
from io import BytesIO
from time import monotonic
from typing import Annotated, Any, Callable, Iterable, Literal
from urllib.parse import parse_qsl, urlparse

import httpx
from docx import Document
from fastapi import APIRouter, Depends, Header, HTTPException, Query, Response, status
from openpyxl import Workbook
from pptx import Presentation
from pydantic import BaseModel, Field
from fastapi.responses import StreamingResponse

from workama_platform.core import Actor, capability_allows, get_actor, json_dumps, new_id, pool, settings
from workama_platform.modules.jobs import IdempotencyConflict, canonical_hash, submit_operation
from workama_platform.object_store import get_object, put_object


router = APIRouter(prefix="/api/v1/work", tags=["ama-work"])
ARTIFACT_BUCKET = "workama-artifacts"
_SENSITIVE_KEY = re.compile(
    r"(?:api[_-]?key|access[_-]?token|refresh[_-]?token|authorization|credential|password|secret|private[_-]?key|token)",
    re.IGNORECASE,
)

PlanStatus = Literal["draft", "ready", "running", "paused", "succeeded", "failed", "cancelled"]
TaskStatus = Literal["todo", "in_progress", "blocked", "done", "failed", "cancelled"]
OfficeFormat = Literal["docx", "xlsx", "pptx"]
ExecutionMode = Literal["requested", "dry_run", "deep_research"]

PLAN_TRANSITIONS: dict[str, frozenset[str]] = {
    "draft": frozenset({"ready", "running", "cancelled"}),
    "ready": frozenset({"running", "cancelled"}),
    "running": frozenset({"paused", "succeeded", "failed", "cancelled"}),
    "paused": frozenset({"running", "cancelled"}),
    "failed": frozenset({"running", "cancelled"}),
    "succeeded": frozenset(),
    "cancelled": frozenset(),
}
TASK_TRANSITIONS: dict[str, frozenset[str]] = {
    "todo": frozenset({"in_progress", "cancelled"}),
    "in_progress": frozenset({"blocked", "done", "failed", "cancelled"}),
    "blocked": frozenset({"in_progress", "cancelled"}),
    "failed": frozenset({"in_progress", "cancelled"}),
    "done": frozenset(),
    "cancelled": frozenset(),
}


class PlanCreate(BaseModel):
    title: str = Field(min_length=1, max_length=160)
    objective: str = Field(default="", max_length=20000)
    session_id: str | None = Field(default=None, max_length=80)


class PlanStatusUpdate(BaseModel):
    status: PlanStatus
    reason: str | None = Field(default=None, max_length=1000)


class TaskCreate(BaseModel):
    title: str = Field(min_length=1, max_length=240)
    description: str = Field(default="", max_length=20000)
    position: int | None = Field(default=None, ge=0, le=100000)


class TaskStatusUpdate(BaseModel):
    status: TaskStatus
    reason: str | None = Field(default=None, max_length=1000)


class TaskReorder(BaseModel):
    task_ids: list[str] = Field(min_length=1, max_length=1000)


class ExecutionRequest(BaseModel):
    mode: ExecutionMode = "requested"
    source_ids: list[str] = Field(default_factory=list, max_length=100)


class ResearchSourceCreate(BaseModel):
    url: str = Field(min_length=1, max_length=2048)
    title: str | None = Field(default=None, max_length=240)
    excerpt: str | None = Field(default=None, max_length=4000)
    fetch: bool = False


class OfficeArtifactCreate(BaseModel):
    format: OfficeFormat
    filename: str | None = Field(default=None, max_length=160)
    upload: bool = True


SCHEMA_STATEMENTS: tuple[str, ...] = (
    """
    CREATE TABLE IF NOT EXISTS work_plan (
        id TEXT PRIMARY KEY,
        workspace_id TEXT NOT NULL REFERENCES id_workspace(id) ON DELETE CASCADE,
        session_id TEXT REFERENCES ag_session(id) ON DELETE SET NULL,
        title TEXT NOT NULL,
        objective TEXT NOT NULL DEFAULT '',
        status TEXT NOT NULL DEFAULT 'draft'
            CHECK (status IN ('draft', 'ready', 'running', 'paused', 'succeeded', 'failed', 'cancelled')),
        last_event_seq BIGINT NOT NULL DEFAULT 0,
        created_by TEXT NOT NULL REFERENCES id_user(id),
        created_at TIMESTAMPTZ NOT NULL DEFAULT now(),
        updated_at TIMESTAMPTZ NOT NULL DEFAULT now()
    )
    """,
    """
    CREATE TABLE IF NOT EXISTS work_execution (
        id TEXT PRIMARY KEY,
        workspace_id TEXT NOT NULL REFERENCES id_workspace(id) ON DELETE CASCADE,
        plan_id TEXT NOT NULL REFERENCES work_plan(id) ON DELETE CASCADE,
        operation_id TEXT NOT NULL UNIQUE REFERENCES ops_async_operation(id) ON DELETE CASCADE,
        source_ids JSONB NOT NULL DEFAULT '[]'::jsonb,
        execution_mode TEXT NOT NULL DEFAULT 'plan'
            CHECK (execution_mode IN ('plan', 'deep_research')),
        status TEXT NOT NULL DEFAULT 'queued'
            CHECK (status IN ('queued', 'running', 'succeeded', 'failed', 'cancelled')),
        requested_by TEXT NOT NULL REFERENCES id_user(id),
        started_at TIMESTAMPTZ,
        completed_at TIMESTAMPTZ,
        created_at TIMESTAMPTZ NOT NULL DEFAULT now(),
        updated_at TIMESTAMPTZ NOT NULL DEFAULT now()
    )
    """,
    """
    CREATE TABLE IF NOT EXISTS work_task (
        id TEXT PRIMARY KEY,
        workspace_id TEXT NOT NULL REFERENCES id_workspace(id) ON DELETE CASCADE,
        plan_id TEXT NOT NULL REFERENCES work_plan(id) ON DELETE CASCADE,
        title TEXT NOT NULL,
        description TEXT NOT NULL DEFAULT '',
        position INTEGER NOT NULL CHECK (position >= 0),
        status TEXT NOT NULL DEFAULT 'todo'
            CHECK (status IN ('todo', 'in_progress', 'blocked', 'done', 'failed', 'cancelled')),
        created_by TEXT NOT NULL REFERENCES id_user(id),
        created_at TIMESTAMPTZ NOT NULL DEFAULT now(),
        updated_at TIMESTAMPTZ NOT NULL DEFAULT now(),
        UNIQUE(plan_id, position)
    )
    """,
    """
    CREATE TABLE IF NOT EXISTS work_event (
        id TEXT PRIMARY KEY,
        workspace_id TEXT NOT NULL REFERENCES id_workspace(id) ON DELETE CASCADE,
        plan_id TEXT NOT NULL REFERENCES work_plan(id) ON DELETE CASCADE,
        task_id TEXT REFERENCES work_task(id) ON DELETE SET NULL,
        seq BIGINT NOT NULL,
        event_type TEXT NOT NULL,
        payload JSONB NOT NULL DEFAULT '{}'::jsonb,
        created_by TEXT NOT NULL REFERENCES id_user(id),
        created_at TIMESTAMPTZ NOT NULL DEFAULT now(),
        UNIQUE(plan_id, seq)
    )
    """,
    """
    CREATE TABLE IF NOT EXISTS work_citation (
        id TEXT PRIMARY KEY,
        workspace_id TEXT NOT NULL REFERENCES id_workspace(id) ON DELETE CASCADE,
        plan_id TEXT NOT NULL REFERENCES work_plan(id) ON DELETE CASCADE,
        task_id TEXT REFERENCES work_task(id) ON DELETE SET NULL,
        source_type TEXT NOT NULL CHECK (source_type IN ('https', 'mock')),
        url TEXT NOT NULL,
        title TEXT,
        excerpt TEXT,
        content_sha256 TEXT,
        created_by TEXT NOT NULL REFERENCES id_user(id),
        created_at TIMESTAMPTZ NOT NULL DEFAULT now(),
        UNIQUE(plan_id, url)
    )
    """,
    """
    CREATE TABLE IF NOT EXISTS work_artifact (
        id TEXT PRIMARY KEY,
        workspace_id TEXT NOT NULL REFERENCES id_workspace(id) ON DELETE CASCADE,
        plan_id TEXT NOT NULL REFERENCES work_plan(id) ON DELETE CASCADE,
        task_id TEXT REFERENCES work_task(id) ON DELETE SET NULL,
        artifact_id TEXT REFERENCES ag_artifact(id) ON DELETE SET NULL,
        name TEXT NOT NULL,
        kind TEXT NOT NULL DEFAULT 'office',
        content_type TEXT NOT NULL,
        s3_key TEXT NOT NULL,
        size_bytes BIGINT NOT NULL CHECK (size_bytes >= 0),
        content_sha256 TEXT NOT NULL,
        status TEXT NOT NULL DEFAULT 'ready' CHECK (status IN ('pending', 'ready', 'failed')),
        preview JSONB NOT NULL DEFAULT '{}'::jsonb,
        created_by TEXT NOT NULL REFERENCES id_user(id),
        created_at TIMESTAMPTZ NOT NULL DEFAULT now()
    )
    """,
    "CREATE INDEX IF NOT EXISTS idx_work_plan_workspace_time ON work_plan(workspace_id, updated_at DESC)",
    "CREATE INDEX IF NOT EXISTS idx_work_execution_plan_time ON work_execution(workspace_id, plan_id, created_at DESC)",
    "CREATE INDEX IF NOT EXISTS idx_work_execution_operation ON work_execution(operation_id)",
    "CREATE INDEX IF NOT EXISTS idx_work_task_plan_position ON work_task(workspace_id, plan_id, position)",
    "CREATE INDEX IF NOT EXISTS idx_work_event_plan_seq ON work_event(workspace_id, plan_id, seq)",
    "CREATE INDEX IF NOT EXISTS idx_work_citation_plan_time ON work_citation(workspace_id, plan_id, created_at DESC)",
    "CREATE INDEX IF NOT EXISTS idx_work_artifact_plan_time ON work_artifact(workspace_id, plan_id, created_at DESC)",
)


async def ensure_work_schema(conn) -> None:
    """Apply the additive AMA-Work schema to an existing connection."""
    for statement in SCHEMA_STATEMENTS:
        await conn.execute(statement)


def validate_plan_transition(current: str, target: str) -> None:
    if target not in PLAN_TRANSITIONS.get(current, frozenset()):
        raise HTTPException(status_code=409, detail=f"Work plan cannot transition from {current} to {target}")


def validate_task_transition(current: str, target: str) -> None:
    if target not in TASK_TRANSITIONS.get(current, frozenset()):
        raise HTTPException(status_code=409, detail=f"Work task cannot transition from {current} to {target}")


def normalize_task_order(existing_ids: Iterable[str], requested_ids: Iterable[str]) -> list[str]:
    existing = list(existing_ids)
    requested = list(requested_ids)
    if len(requested) != len(set(requested)):
        raise HTTPException(status_code=422, detail="Task order contains duplicate IDs")
    if set(existing) != set(requested) or len(existing) != len(requested):
        raise HTTPException(status_code=422, detail="Task order must contain exactly the plan tasks")
    return requested


def next_task_position(tasks: Iterable[dict[str, Any]]) -> int:
    """Return a free position for the temporary insert before reordering."""
    return max((int(task.get("position", -1)) for task in tasks), default=-1) + 1


def redact_sensitive(value: Any) -> Any:
    if isinstance(value, dict):
        return {
            str(key): "<redacted>" if _SENSITIVE_KEY.search(str(key)) else redact_sensitive(item)
            for key, item in value.items()
        }
    if isinstance(value, list):
        return [redact_sensitive(item) for item in value]
    if isinstance(value, tuple):
        return [redact_sensitive(item) for item in value]
    return value


def _blocked_hostname(hostname: str) -> bool:
    normalized = hostname.rstrip(".").lower()
    return normalized in {"localhost", "localhost.localdomain", "0.0.0.0", "::1"} or normalized.endswith(
        (".localhost", ".local", ".internal")
    )


def _blocked_ip(value: str) -> bool:
    try:
        address = ipaddress.ip_address(value)
    except ValueError:
        return False
    return not address.is_global or address.is_private or address.is_loopback or address.is_link_local


def validate_research_url(
    value: str,
    *,
    resolver: Callable[..., list[tuple[Any, ...]]] | None = None,
    resolve_dns: bool = True,
) -> dict[str, str]:
    parsed = urlparse(value.strip())
    if parsed.scheme not in {"https", "mock"} or not parsed.hostname:
        raise HTTPException(status_code=422, detail="Research sources must use https:// or mock://")
    if parsed.username or parsed.password:
        raise HTTPException(status_code=422, detail="Research source credentials are not allowed")
    if _blocked_hostname(parsed.hostname) or _blocked_ip(parsed.hostname):
        raise HTTPException(status_code=422, detail="Research source host is not allowed")
    try:
        port = parsed.port or 443
    except ValueError as exc:
        raise HTTPException(status_code=422, detail="Research source port is invalid") from exc
    for key, _ in parse_qsl(parsed.query, keep_blank_values=True):
        if _SENSITIVE_KEY.search(key):
            raise HTTPException(status_code=422, detail="Research source credential query parameters are not allowed")
    # 运行时查找 socket.getaddrinfo，确保 monkeypatch 替换模块属性后生效
    if resolver is None:
        resolver = socket.getaddrinfo
    if parsed.scheme == "https" and resolve_dns:
        try:
            addresses = resolver(parsed.hostname, port, type=socket.SOCK_STREAM)
        except (OSError, socket.gaierror) as exc:
            raise HTTPException(status_code=422, detail="Research source host cannot be resolved safely") from exc
        for address in addresses:
            if len(address) > 4 and _blocked_ip(address[4][0]):
                raise HTTPException(status_code=422, detail="Research source resolves to a private address")
    return {
        "url": parsed.geturl(),
        "source_type": parsed.scheme,
        "host": parsed.hostname,
    }


def deterministic_mock_browser_fetch(url: str) -> dict[str, Any]:
    validated = validate_research_url(url, resolve_dns=False)
    if validated["source_type"] != "mock":
        raise HTTPException(status_code=422, detail="Deterministic browser fetch only accepts mock:// sources")
    text = f"Deterministic WorkAMA research result for {validated['url']}."
    return {
        "url": validated["url"],
        "title": f"Mock source: {validated['host']}",
        "text": text,
        "content_sha256": hashlib.sha256(text.encode()).hexdigest(),
        "untrusted": True,
    }


async def sandbox_browser_fetch(url: str, *, client: httpx.AsyncClient | None = None) -> dict[str, Any]:
    """通过 sandbox-browser 容器真实抓取 https 研究源。

    流程：申请 sandbox-browser 沙箱 → navigate → eval(document.title) →
    eval(document.body.innerText.slice(0,5000)) → screenshot → 释放沙箱。
    任何步骤失败抛 HTTPException(503)；finally 块确保沙箱释放。
    抓取结果统一标记 ``untrusted=True``：浏览器自动抓取的内容不可作为可信证据。
    """
    validated = validate_research_url(url)
    if validated["source_type"] != "https":
        raise HTTPException(status_code=422, detail="Sandbox browser fetch only accepts https:// sources")

    fleet_base = settings.sandbox_fleet_url.rstrip("/")
    headers = {"X-Internal-Token": settings.internal_token}
    # session_id 由 URL 派生，保证同一 URL 复用同一沙箱 scope
    session_id = hashlib.sha256(validated["url"].encode()).hexdigest()[:16]
    sandbox_id: str | None = None
    owns_client = client is None
    if client is None:
        client = httpx.AsyncClient(timeout=30.0, follow_redirects=False, trust_env=False)

    try:
        # 1. 申请 sandbox-browser 沙箱
        acquire_resp = await client.post(
            f"{fleet_base}/internal/sandboxes",
            headers=headers,
            json={
                "workspace_id": "ama-work-research",
                "session_id": session_id,
                "image": "sandbox-browser",
            },
        )
        if acquire_resp.status_code >= 400:
            raise HTTPException(
                status_code=503,
                detail=f"Sandbox browser fetch failed: acquire returned {acquire_resp.status_code}",
            )
        sandbox_id = str(acquire_resp.json().get("id") or "")
        if not sandbox_id:
            raise HTTPException(status_code=503, detail="Sandbox browser fetch failed: no sandbox id returned")

        # 2. navigate 到目标 URL
        navigate_resp = await client.post(
            f"{fleet_base}/internal/sandboxes/{sandbox_id}/browser",
            headers=headers,
            json={"action": "navigate", "target": validated["url"], "timeout_ms": 20000},
        )
        if navigate_resp.status_code >= 400:
            raise HTTPException(
                status_code=503,
                detail=f"Sandbox browser fetch failed: navigate returned {navigate_resp.status_code}",
            )
        navigate_payload = navigate_resp.json()
        if not navigate_payload.get("ok"):
            raise HTTPException(
                status_code=503,
                detail=f"Sandbox browser fetch failed: navigate error={navigate_payload.get('error')}",
            )

        # 3. eval(document.title) 取标题
        title_resp = await client.post(
            f"{fleet_base}/internal/sandboxes/{sandbox_id}/browser",
            headers=headers,
            json={
                "action": "eval",
                "params": {"expression": "document.title"},
                "timeout_ms": 10000,
            },
        )
        if title_resp.status_code >= 400 or not title_resp.json().get("ok"):
            raise HTTPException(
                status_code=503,
                detail="Sandbox browser fetch failed: eval document.title failed",
            )
        title = str((title_resp.json().get("meta") or {}).get("result") or "")

        # 4. eval(document.body.innerText.slice(0, 5000)) 取正文摘录
        text_resp = await client.post(
            f"{fleet_base}/internal/sandboxes/{sandbox_id}/browser",
            headers=headers,
            json={
                "action": "eval",
                "params": {"expression": "document.body.innerText.slice(0, 5000)"},
                "timeout_ms": 10000,
            },
        )
        if text_resp.status_code >= 400 or not text_resp.json().get("ok"):
            raise HTTPException(
                status_code=503,
                detail="Sandbox browser fetch failed: eval body.innerText failed",
            )
        text = str((text_resp.json().get("meta") or {}).get("result") or "")

        # 5. screenshot 取截图（base64 PNG）
        screenshot_resp = await client.post(
            f"{fleet_base}/internal/sandboxes/{sandbox_id}/browser",
            headers=headers,
            json={"action": "screenshot", "timeout_ms": 10000},
        )
        if screenshot_resp.status_code >= 400 or not screenshot_resp.json().get("ok"):
            raise HTTPException(
                status_code=503,
                detail="Sandbox browser fetch failed: screenshot failed",
            )
        screenshot = str(screenshot_resp.json().get("screenshot") or "")

        return {
            "url": validated["url"],
            "title": title,
            "text": text,
            "content_sha256": hashlib.sha256(text.encode()).hexdigest(),
            "screenshot": screenshot,
            "untrusted": True,
        }
    except HTTPException:
        raise
    except httpx.HTTPError as exc:
        raise HTTPException(status_code=503, detail=f"Sandbox browser fetch failed: {exc}") from exc
    finally:
        # 释放沙箱：无论成功或失败，都尝试 DELETE，错误被吞掉避免掩盖原始异常
        if sandbox_id:
            try:
                await client.delete(
                    f"{fleet_base}/internal/sandboxes/{sandbox_id}",
                    headers=headers,
                )
            except httpx.HTTPError:
                pass
        if owns_client:
            await client.aclose()


@dataclass(frozen=True)
class OfficeArtifact:
    data: bytes
    extension: str
    content_type: str


def research_cross_validation(records: Iterable[dict[str, Any]]) -> dict[str, Any]:
    """Return conservative, fingerprint-level validation for controlled sources."""
    rows = list(records)
    fingerprints = sorted(
        {
            str((record.get("fetched") or {}).get("content_sha256") or "")
            for record in rows
            if (record.get("fetched") or {}).get("content_sha256")
        }
    )
    return {
        "source_count": len(rows),
        "distinct_fingerprint_count": len(fingerprints),
        "status": "consistent_fingerprint" if len(fingerprints) <= 1 else "divergent_fingerprint",
        "fingerprints": fingerprints,
        "trust_boundary": "mock_sources_are_untrusted",
    }


def generate_research_markdown(
    title: str,
    objective: str,
    records: Iterable[dict[str, Any]],
    validation: dict[str, Any],
) -> str:
    rows = list(records)
    lines = [
        f"# {redact_sensitive(title).strip() or 'WorkAMA deep research report'}",
        "",
        "> Local controlled evidence report. `mock://` sources are untrusted fixtures and are not proof of external facts.",
        "",
        "## Objective",
        "",
        redact_sensitive(objective).strip() or "No objective was provided.",
        "",
        "## Research rounds",
        "",
        "1. Source collection: deterministic controlled browser fixtures.",
        "2. Cross-validation: compare source content fingerprints and preserve disagreements.",
        "",
        "## Findings",
        "",
    ]
    if rows:
        for index, record in enumerate(rows, start=1):
            source = record.get("source") or {}
            fetched = record.get("fetched") or {}
            text = re.sub(r"\s+", " ", str(fetched.get("text") or source.get("excerpt") or "No source excerpt available.")).strip()
            lines.append(f"{index}. {text[:800]} [{index}]")
    else:
        lines.append("No controlled sources were available.")
    lines.extend(
        [
            "",
            "## Cross-validation",
            "",
            f"- Sources reviewed: {validation.get('source_count', len(rows))}",
            f"- Distinct content fingerprints: {validation.get('distinct_fingerprint_count', 0)}",
            f"- Result: `{validation.get('status', 'unknown')}`",
            "- This check compares captured content hashes only; it does not establish factual truth or source independence.",
            "",
            "## References",
            "",
        ]
    )
    if rows:
        for index, record in enumerate(rows, start=1):
            source = record.get("source") or {}
            fetched = record.get("fetched") or {}
            lines.append(
                f"[{index}] {source.get('title') or fetched.get('title') or source.get('url') or 'Untitled source'} - "
                f"{source.get('url', '')} - SHA-256: {fetched.get('content_sha256') or source.get('content_sha256') or 'not captured'}"
            )
    else:
        lines.append("No references.")
    return "\n".join(lines) + "\n"


def _pdf_escape(value: str) -> str:
    ascii_value = str(value).encode("ascii", "replace").decode("ascii")
    return ascii_value.replace("\\", "\\\\").replace("(", "\\(").replace(")", "\\)")


def generate_simple_pdf(markdown: str) -> bytes:
    """Create a dependency-free, text PDF for deterministic local report evidence."""
    lines: list[str] = []
    for raw in markdown.splitlines():
        text = raw.replace("`", "").replace("**", "")
        if not text:
            lines.append("")
            continue
        while len(text) > 96:
            lines.append(text[:96])
            text = text[96:]
        lines.append(text)
    lines = lines[:48]
    stream_lines = ["BT", "/F1 10 Tf", "50 760 Td", "14 TL"]
    stream_lines.extend(f"({_pdf_escape(line)}) Tj T*" for line in lines)
    stream_lines.append("ET")
    stream = "\n".join(stream_lines).encode("ascii")
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Resources << /Font << /F1 4 0 R >> >> /Contents 5 0 R >>",
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
        b"<< /Length " + str(len(stream)).encode("ascii") + b" >>\nstream\n" + stream + b"\nendstream",
    ]
    output = bytearray(b"%PDF-1.4\n%\xe2\xe3\xcf\xd3\n")
    offsets = [0]
    for index, obj in enumerate(objects, start=1):
        offsets.append(len(output))
        output.extend(f"{index} 0 obj\n".encode("ascii"))
        output.extend(obj)
        output.extend(b"\nendobj\n")
    xref_offset = len(output)
    output.extend(f"xref\n0 {len(objects) + 1}\n".encode("ascii"))
    output.extend(b"0000000000 65535 f \n")
    for offset in offsets[1:]:
        output.extend(f"{offset:010d} 00000 n \n".encode("ascii"))
    output.extend(
        f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\nstartxref\n{xref_offset}\n%%EOF\n".encode("ascii")
    )
    return bytes(output)


def generate_research_artifacts(
    title: str,
    objective: str,
    records: Iterable[dict[str, Any]],
) -> tuple[OfficeArtifact, OfficeArtifact, dict[str, Any]]:
    rows = list(records)
    validation = research_cross_validation(rows)
    markdown = generate_research_markdown(title, objective, rows, validation)
    return (
        OfficeArtifact(markdown.encode("utf-8"), "md", "text/markdown"),
        OfficeArtifact(generate_simple_pdf(markdown), "pdf", "application/pdf"),
        validation,
    )


def _task_rows_for_office(tasks: Iterable[dict[str, Any]]) -> list[dict[str, str]]:
    return [
        {"title": str(task.get("title", "")), "status": str(task.get("status", "todo")), "description": str(task.get("description", ""))}
        for task in tasks
    ]


def generate_docx(title: str, objective: str, tasks: Iterable[dict[str, Any]], sources: Iterable[dict[str, Any]]) -> OfficeArtifact:
    document = Document()
    document.add_heading(title, level=0)
    if objective:
        document.add_paragraph(objective)
    document.add_heading("Task list", level=1)
    for task in _task_rows_for_office(tasks):
        document.add_paragraph(f"[{task['status']}] {task['title']}", style="List Bullet")
        if task["description"]:
            document.add_paragraph(task["description"])
    source_rows = list(sources)
    if source_rows:
        document.add_heading("Sources", level=1)
        for source in source_rows:
            document.add_paragraph(str(source.get("url", "")), style="List Bullet")
    output = BytesIO()
    document.save(output)
    return OfficeArtifact(output.getvalue(), "docx", "application/vnd.openxmlformats-officedocument.wordprocessingml.document")


def generate_xlsx(title: str, objective: str, tasks: Iterable[dict[str, Any]], sources: Iterable[dict[str, Any]]) -> OfficeArtifact:
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Plan"
    sheet.append([title])
    sheet.append([objective])
    sheet.append([])
    sheet.append(["Position", "Task", "Status", "Description"])
    for position, task in enumerate(_task_rows_for_office(tasks), start=1):
        sheet.append([position, task["title"], task["status"], task["description"]])
    source_sheet = workbook.create_sheet("Sources")
    source_sheet.append(["URL", "Title", "SHA256"])
    for source in sources:
        source_sheet.append([source.get("url", ""), source.get("title", ""), source.get("content_sha256", "")])
    output = BytesIO()
    workbook.save(output)
    return OfficeArtifact(output.getvalue(), "xlsx", "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")


def generate_pptx(title: str, objective: str, tasks: Iterable[dict[str, Any]], sources: Iterable[dict[str, Any]]) -> OfficeArtifact:
    source_rows = list(sources)
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = title
    frame = slide.placeholders[1].text_frame
    frame.text = objective or "AMA-Work plan"
    for task in _task_rows_for_office(tasks):
        paragraph = frame.add_paragraph()
        paragraph.text = f"[{task['status']}] {task['title']}"
        paragraph.level = 0
    if source_rows:
        source_slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        source_slide.shapes.title.text = "Sources"
        source_frame = source_slide.placeholders[1].text_frame
        source_frame.text = ""
        for source in source_rows:
            paragraph = source_frame.add_paragraph()
            paragraph.text = str(source.get("url", ""))
    output = BytesIO()
    presentation.save(output)
    return OfficeArtifact(output.getvalue(), "pptx", "application/vnd.openxmlformats-officedocument.presentationml.presentation")


def generate_office_artifact(
    format: OfficeFormat,
    title: str,
    objective: str,
    tasks: Iterable[dict[str, Any]],
    sources: Iterable[dict[str, Any]],
) -> OfficeArtifact:
    task_rows = list(tasks)
    source_rows = list(sources)
    if format == "docx":
        return generate_docx(title, objective, task_rows, source_rows)
    if format == "xlsx":
        return generate_xlsx(title, objective, task_rows, source_rows)
    return generate_pptx(title, objective, task_rows, source_rows)


def _safe_filename(value: str, extension: str) -> str:
    name = re.sub(r"[^A-Za-z0-9._-]+", "-", value).strip(".-")[:120] or "work-plan"
    if not name.lower().endswith(f".{extension}"):
        name = f"{name}.{extension}"
    return name


def _require(actor: Actor, action: str) -> None:
    if actor.actor_type != "user":
        raise HTTPException(status_code=403, detail="User authentication is required")
    if action == "read":
        return
    if capability_allows(actor.capabilities, f"work:{action}"):
        return
    if action == "write" and actor.role in {"owner", "admin", "member"}:
        return
    raise HTTPException(status_code=403, detail=f"Missing capability: work:{action}")


def plan_view(row: dict[str, Any]) -> dict[str, Any]:
    result = {
        "id": row["id"],
        "workspace_id": row["workspace_id"],
        "session_id": row.get("session_id"),
        "title": row["title"],
        "objective": row.get("objective", ""),
        "status": row["status"],
        "last_event_seq": row.get("last_event_seq", 0),
        "created_by": row["created_by"],
        "created_at": row.get("created_at"),
        "updated_at": row.get("updated_at"),
    }
    if "latest_execution" in row:
        result["latest_execution"] = execution_view(row.get("latest_execution")) if row.get("latest_execution") else None
    return result


def execution_view(row: dict[str, Any] | None) -> dict[str, Any] | None:
    if not row:
        return None
    source_ids = row.get("source_ids") or []
    return {
        "id": row["id"],
        "workspace_id": row["workspace_id"],
        "plan_id": row["plan_id"],
        "operation_id": row["operation_id"],
        "source_ids": source_ids if isinstance(source_ids, list) else [],
        "execution_mode": row.get("execution_mode", "plan"),
        "status": row.get("status"),
        "operation_status": row.get("operation_status", row.get("status")),
        "progress": row.get("progress", 0),
        "stage": row.get("stage"),
        "started_at": row.get("started_at"),
        "completed_at": row.get("completed_at"),
        "created_at": row.get("created_at"),
        "updated_at": row.get("updated_at"),
    }


def task_view(row: dict[str, Any]) -> dict[str, Any]:
    return {
        "id": row["id"],
        "workspace_id": row["workspace_id"],
        "plan_id": row["plan_id"],
        "title": row["title"],
        "description": row.get("description", ""),
        "position": row["position"],
        "status": row["status"],
        "created_by": row["created_by"],
        "created_at": row.get("created_at"),
        "updated_at": row.get("updated_at"),
    }


def event_view(row: dict[str, Any]) -> dict[str, Any]:
    return {
        "id": row["id"],
        "workspace_id": row["workspace_id"],
        "plan_id": row["plan_id"],
        "task_id": row.get("task_id"),
        "seq": row["seq"],
        "event_type": row["event_type"],
        "payload": redact_sensitive(row.get("payload") or {}),
        "created_by": row["created_by"],
        "created_at": row.get("created_at"),
    }


async def _owned_plan(conn, plan_id: str, actor: Actor, *, for_update: bool = False) -> dict[str, Any]:
    lock = " FOR UPDATE" if for_update else ""
    result = await conn.execute(
        f"""
        SELECT id, workspace_id, session_id, title, objective, status,
               last_event_seq, created_by, created_at, updated_at
        FROM work_plan
        WHERE id=%s AND workspace_id=%s{lock}
        """,
        (plan_id, actor.workspace_id),
    )
    row = await result.fetchone()
    if not row:
        raise HTTPException(status_code=404, detail="Work plan not found")
    return row


async def _owned_task(conn, plan_id: str, task_id: str, actor: Actor, *, for_update: bool = False) -> dict[str, Any]:
    lock = " FOR UPDATE" if for_update else ""
    result = await conn.execute(
        f"""
        SELECT id, workspace_id, plan_id, title, description, position, status,
               created_by, created_at, updated_at
        FROM work_task
        WHERE id=%s AND plan_id=%s AND workspace_id=%s{lock}
        """,
        (task_id, plan_id, actor.workspace_id),
    )
    row = await result.fetchone()
    if not row:
        raise HTTPException(status_code=404, detail="Work task not found")
    return row


async def _append_event(
    conn,
    *,
    plan: dict[str, Any],
    actor: Actor,
    event_type: str,
    payload: dict[str, Any],
    task_id: str | None = None,
) -> dict[str, Any]:
    safe_payload = redact_sensitive(payload)
    seq = int(plan.get("last_event_seq") or 0) + 1
    await conn.execute(
        "UPDATE work_plan SET last_event_seq=%s, updated_at=now() WHERE id=%s AND workspace_id=%s",
        (seq, plan["id"], actor.workspace_id),
    )
    result = await conn.execute(
        """
        INSERT INTO work_event(
            id, workspace_id, plan_id, task_id, seq, event_type, payload, created_by
        ) VALUES (%s,%s,%s,%s,%s,%s,%s::jsonb,%s)
        RETURNING id, workspace_id, plan_id, task_id, seq, event_type, payload, created_by, created_at
        """,
        (
            new_id("wevt"), actor.workspace_id, plan["id"], task_id, seq,
            event_type, json_dumps(safe_payload), actor.user_id,
        ),
    )
    plan["last_event_seq"] = seq
    return event_view(await result.fetchone())


async def _task_rows(conn, plan_id: str, actor: Actor, *, for_update: bool = False) -> list[dict[str, Any]]:
    lock = " FOR UPDATE" if for_update else ""
    result = await conn.execute(
        f"""
        SELECT id, workspace_id, plan_id, title, description, position, status,
               created_by, created_at, updated_at
        FROM work_task
        WHERE plan_id=%s AND workspace_id=%s
        ORDER BY position, id{lock}
        """,
        (plan_id, actor.workspace_id),
    )
    return await result.fetchall()


async def _persist_order(conn, plan_id: str, workspace_id: str, task_ids: list[str]) -> None:
    if not task_ids:
        return
    offset = 1_000_000
    await conn.execute(
        "UPDATE work_task SET position=position+%s WHERE plan_id=%s AND workspace_id=%s",
        (offset, plan_id, workspace_id),
    )
    for position, task_id in enumerate(task_ids):
        await conn.execute(
            "UPDATE work_task SET position=%s, updated_at=now() WHERE id=%s AND plan_id=%s AND workspace_id=%s",
            (position, task_id, plan_id, workspace_id),
        )


@router.post("/plans", status_code=status.HTTP_201_CREATED)
async def create_plan(body: PlanCreate, actor: Annotated[Actor, Depends(get_actor)]):
    _require(actor, "write")
    plan_id = new_id("wplan")
    async with pool.connection() as conn:
        if body.session_id:
            result = await conn.execute(
                "SELECT id FROM ag_session WHERE id=%s AND workspace_id=%s",
                (body.session_id, actor.workspace_id),
            )
            if not await result.fetchone():
                raise HTTPException(status_code=404, detail="Agent session not found")
        async with conn.transaction():
            await conn.execute(
                """
                INSERT INTO work_plan(id,workspace_id,session_id,title,objective,created_by)
                VALUES (%s,%s,%s,%s,%s,%s)
                """,
                (plan_id, actor.workspace_id, body.session_id, body.title.strip(), body.objective.strip(), actor.user_id),
            )
            plan = await _owned_plan(conn, plan_id, actor, for_update=True)
            await _append_event(conn, plan=plan, actor=actor, event_type="plan.created", payload={"title": plan["title"]})
    return plan_view(plan)


@router.get("/plans")
async def list_plans(
    actor: Annotated[Actor, Depends(get_actor)],
    plan_status: PlanStatus | None = Query(default=None, alias="status"),
    limit: int = Query(default=100, ge=1, le=200),
):
    _require(actor, "read")
    clauses = ["workspace_id=%s"]
    params: list[Any] = [actor.workspace_id]
    if plan_status:
        clauses.append("status=%s")
        params.append(plan_status)
    params.append(limit)
    async with pool.connection() as conn:
        result = await conn.execute(
            f"""
            SELECT id,workspace_id,session_id,title,objective,status,last_event_seq,
                   created_by,created_at,updated_at
            FROM work_plan WHERE {' AND '.join(clauses)}
            ORDER BY updated_at DESC, id DESC LIMIT %s
            """,
            tuple(params),
        )
        rows = await result.fetchall()
    # Contract《720》listPlans: ListQuery -> ListResponse<WorkPlanDTO>
    data = [plan_view(row) for row in rows]
    return {
        "items": data,
        "data": data,
        "next_cursor": None,
        "has_more": False,
        "meta": {"request_id": None, "count": len(data)},
    }


@router.get("/plans/{plan_id}")
async def get_plan(plan_id: str, actor: Annotated[Actor, Depends(get_actor)]):
    _require(actor, "read")
    async with pool.connection() as conn:
        plan = await _owned_plan(conn, plan_id, actor)
        tasks = await _task_rows(conn, plan_id, actor)
        sources = await conn.execute(
            """
            SELECT id,plan_id,task_id,source_type,url,title,excerpt,content_sha256,created_by,created_at
            FROM work_citation WHERE plan_id=%s AND workspace_id=%s ORDER BY created_at,id
            """,
            (plan_id, actor.workspace_id),
        )
        execution_result = await conn.execute(
            """
            SELECT e.id,e.workspace_id,e.plan_id,e.operation_id,e.source_ids,e.execution_mode,e.status,
                   e.started_at,e.completed_at,e.created_at,e.updated_at,
                   o.status AS operation_status,o.progress,o.stage
            FROM work_execution e
            JOIN ops_async_operation o ON o.id=e.operation_id
            WHERE e.plan_id=%s AND e.workspace_id=%s
            ORDER BY e.created_at DESC, e.id DESC
            LIMIT 1
            """,
            (plan_id, actor.workspace_id),
        )
        plan_view_data = plan_view(plan)
        plan_view_data["tasks"] = [task_view(row) for row in tasks]
        plan_view_data["sources"] = await sources.fetchall()
        plan_view_data["latest_execution"] = execution_view(await execution_result.fetchone())
    return plan_view_data


@router.post("/plans/{plan_id}/status")
async def update_plan_status(plan_id: str, body: PlanStatusUpdate, actor: Annotated[Actor, Depends(get_actor)]):
    _require(actor, "write")
    async with pool.connection() as conn:
        async with conn.transaction():
            plan = await _owned_plan(conn, plan_id, actor, for_update=True)
            validate_plan_transition(plan["status"], body.status)
            await conn.execute(
                "UPDATE work_plan SET status=%s,updated_at=now() WHERE id=%s AND workspace_id=%s",
                (body.status, plan_id, actor.workspace_id),
            )
            event = await _append_event(
                conn, plan=plan, actor=actor, event_type="plan.status.updated",
                payload={"previous_status": plan["status"], "status": body.status, "reason": body.reason},
            )
            plan["status"] = body.status
    return {"previous_status": event["payload"].get("previous_status"), "plan": plan_view(plan), "event": event}


@router.post("/plans/{plan_id}/tasks", status_code=status.HTTP_201_CREATED)
async def create_task(plan_id: str, body: TaskCreate, actor: Annotated[Actor, Depends(get_actor)]):
    _require(actor, "write")
    task_id = new_id("wtask")
    async with pool.connection() as conn:
        async with conn.transaction():
            plan = await _owned_plan(conn, plan_id, actor, for_update=True)
            tasks = await _task_rows(conn, plan_id, actor, for_update=True)
            position = len(tasks) if body.position is None else min(body.position, len(tasks))
            temporary_position = next_task_position(tasks)
            ordered_ids = [row["id"] for row in tasks]
            ordered_ids.insert(position, task_id)
            await conn.execute(
                """
                INSERT INTO work_task(id,workspace_id,plan_id,title,description,position,created_by)
                VALUES (%s,%s,%s,%s,%s,%s,%s)
                """,
                (task_id, actor.workspace_id, plan_id, body.title.strip(), body.description.strip(), temporary_position, actor.user_id),
            )
            await _persist_order(conn, plan_id, actor.workspace_id, ordered_ids)
            result = await conn.execute(
                """
                SELECT id,workspace_id,plan_id,title,description,position,status,created_by,created_at,updated_at
                FROM work_task WHERE id=%s AND plan_id=%s AND workspace_id=%s
                """,
                (task_id, plan_id, actor.workspace_id),
            )
            task = await result.fetchone()
            await _append_event(conn, plan=plan, actor=actor, task_id=task_id, event_type="task.created", payload={"position": position, "title": task["title"]})
    return task_view(task)


@router.get("/plans/{plan_id}/tasks")
async def list_tasks(plan_id: str, actor: Annotated[Actor, Depends(get_actor)]):
    _require(actor, "read")
    async with pool.connection() as conn:
        await _owned_plan(conn, plan_id, actor)
        rows = await _task_rows(conn, plan_id, actor)
    # Contract《720》listPlanTasks: ListQuery -> ListResponse<WorkTaskDTO>
    data = [task_view(row) for row in rows]
    return {
        "items": data,
        "data": data,
        "next_cursor": None,
        "has_more": False,
        "meta": {"request_id": None, "count": len(data)},
    }


@router.delete("/plans/{plan_id}/tasks/{task_id}", status_code=status.HTTP_204_NO_CONTENT)
async def delete_task(plan_id: str, task_id: str, actor: Annotated[Actor, Depends(get_actor)]):
    _require(actor, "write")
    async with pool.connection() as conn:
        async with conn.transaction():
            plan = await _owned_plan(conn, plan_id, actor, for_update=True)
            await _owned_task(conn, plan_id, task_id, actor, for_update=True)
            event = await _append_event(conn, plan=plan, actor=actor, task_id=task_id, event_type="task.deleted", payload={"task_id": task_id})
            await conn.execute(
                "DELETE FROM work_task WHERE id=%s AND plan_id=%s AND workspace_id=%s",
                (task_id, plan_id, actor.workspace_id),
            )
            remaining = await _task_rows(conn, plan_id, actor, for_update=True)
            await _persist_order(conn, plan_id, actor.workspace_id, [row["id"] for row in remaining])
    return Response(status_code=status.HTTP_204_NO_CONTENT)


@router.post("/plans/{plan_id}/tasks/reorder")
async def reorder_tasks(plan_id: str, body: TaskReorder, actor: Annotated[Actor, Depends(get_actor)]):
    _require(actor, "write")
    async with pool.connection() as conn:
        async with conn.transaction():
            plan = await _owned_plan(conn, plan_id, actor, for_update=True)
            tasks = await _task_rows(conn, plan_id, actor, for_update=True)
            ordered_ids = normalize_task_order([row["id"] for row in tasks], body.task_ids)
            await _persist_order(conn, plan_id, actor.workspace_id, ordered_ids)
            event = await _append_event(conn, plan=plan, actor=actor, event_type="task.reordered", payload={"task_ids": ordered_ids})
    return {"task_ids": ordered_ids, "event": event}


@router.post("/plans/{plan_id}/tasks/{task_id}/status")
async def update_task_status(plan_id: str, task_id: str, body: TaskStatusUpdate, actor: Annotated[Actor, Depends(get_actor)]):
    _require(actor, "write")
    async with pool.connection() as conn:
        async with conn.transaction():
            plan = await _owned_plan(conn, plan_id, actor, for_update=True)
            task = await _owned_task(conn, plan_id, task_id, actor, for_update=True)
            validate_task_transition(task["status"], body.status)
            await conn.execute(
                "UPDATE work_task SET status=%s,updated_at=now() WHERE id=%s AND plan_id=%s AND workspace_id=%s",
                (body.status, task_id, plan_id, actor.workspace_id),
            )
            event = await _append_event(
                conn, plan=plan, actor=actor, task_id=task_id, event_type="task.status.updated",
                payload={"previous_status": task["status"], "status": body.status, "reason": body.reason},
            )
            task["status"] = body.status
    return {"previous_status": event["payload"].get("previous_status"), "task": task_view(task), "event": event}


@router.post("/plans/{plan_id}/executions", status_code=status.HTTP_202_ACCEPTED)
async def request_execution(
    plan_id: str,
    body: ExecutionRequest,
    actor: Annotated[Actor, Depends(get_actor)],
    idempotency_key: Annotated[str | None, Header(alias="Idempotency-Key", min_length=1, max_length=200)] = None,
):
    _require(actor, "write")
    async with pool.connection() as conn:
        async with conn.transaction():
            plan = await _owned_plan(conn, plan_id, actor, for_update=True)
            if plan["status"] in {"succeeded", "cancelled"}:
                raise HTTPException(status_code=409, detail="Terminal work plan cannot be executed")
            if len(body.source_ids) != len(set(body.source_ids)):
                raise HTTPException(status_code=422, detail="Execution source_ids must be unique")
            if body.source_ids:
                source_result = await conn.execute(
                    "SELECT id FROM work_citation WHERE plan_id=%s AND workspace_id=%s AND id=ANY(%s)",
                    (plan_id, actor.workspace_id, body.source_ids),
                )
                found = {row["id"] for row in await source_result.fetchall()}
                if found != set(body.source_ids):
                    raise HTTPException(status_code=404, detail="One or more work sources not found")
            if body.mode == "dry_run":
                event = await _append_event(
                    conn, plan=plan, actor=actor, event_type="plan.execution.requested",
                    payload={"mode": body.mode, "source_ids": body.source_ids, "previous_status": plan["status"]},
                )
                return {"plan_id": plan_id, "status": plan["status"], "event": event}

            request_key = idempotency_key or new_id("work-idem")
            stable_input_hash = canonical_hash({
                "plan_id": plan_id,
                "source_ids": body.source_ids,
                "mode": body.mode,
            })
            payload = {
                "plan_id": plan_id,
                "plan_title": plan["title"],
                "plan_objective": plan.get("objective", ""),
                "source_ids": body.source_ids,
                "execution_mode": "deep_research" if body.mode == "deep_research" else "plan",
                "actor_id": actor.user_id,
                "actor_role": actor.role,
                "org_id": actor.org_id,
            }
            try:
                operation = await submit_operation(
                    conn,
                    operation_type="work.plan.execute",
                    workspace_id=actor.workspace_id,
                    org_id=actor.org_id,
                    actor_id=actor.user_id,
                    actor_role=actor.role,
                    idempotency_key=request_key,
                    payload=payload,
                    input_hash_override=stable_input_hash,
                    job_type="work.plan.execute",
                    queue="platform",
                    max_attempts=1,
                    priority=115,
                    cancellable=True,
                )
            except IdempotencyConflict as exc:
                raise HTTPException(status_code=409, detail="Idempotency key was already used with different input") from exc

            existing_execution_result = await conn.execute(
                """
                SELECT e.id,e.workspace_id,e.plan_id,e.operation_id,e.source_ids,e.execution_mode,e.status,
                       e.started_at,e.completed_at,e.created_at,e.updated_at,
                       o.status AS operation_status,o.progress,o.stage
                FROM work_execution e
                JOIN ops_async_operation o ON o.id=e.operation_id
                WHERE e.operation_id=%s AND e.workspace_id=%s
                """,
                (operation["id"], actor.workspace_id),
            )
            existing_execution = await existing_execution_result.fetchone()
            if existing_execution:
                return {
                    "plan_id": plan_id,
                    "status": plan["status"],
                    "operation_id": operation["id"],
                    "execution": execution_view(existing_execution),
                    "event": None,
                }

            execution_id = new_id("wexec")
            previous_status = plan["status"]
            if previous_status != "running":
                validate_plan_transition(previous_status, "running")
                await conn.execute("UPDATE work_plan SET status='running',updated_at=now() WHERE id=%s AND workspace_id=%s", (plan_id, actor.workspace_id))
                plan["status"] = "running"
            await conn.execute(
                """
                INSERT INTO work_execution(
                    id,workspace_id,plan_id,operation_id,source_ids,execution_mode,status,requested_by
                ) VALUES (%s,%s,%s,%s,%s::jsonb,%s,'queued',%s)
                """,
                (
                    execution_id,
                    actor.workspace_id,
                    plan_id,
                    operation["id"],
                    json_dumps(body.source_ids),
                    "deep_research" if body.mode == "deep_research" else "plan",
                    actor.user_id,
                ),
            )
            event = await _append_event(
                conn, plan=plan, actor=actor, event_type="plan.execution.requested",
                payload={
                    "mode": body.mode,
                    "source_ids": body.source_ids,
                    "previous_status": previous_status,
                    "operation_id": operation["id"],
                    "execution_id": execution_id,
                },
            )
    return {
        "plan_id": plan_id,
        "status": plan["status"],
        "operation_id": operation["id"],
        "execution_id": execution_id,
        "execution_status": operation["status"],
        "event": event,
    }


@router.get("/plans/{plan_id}/events")
async def list_events(
    plan_id: str,
    actor: Annotated[Actor, Depends(get_actor)],
    after: int = Query(default=0, ge=0),
    limit: int = Query(default=500, ge=1, le=1000),
):
    _require(actor, "read")
    async with pool.connection() as conn:
        await _owned_plan(conn, plan_id, actor)
        result = await conn.execute(
            """
            SELECT id,workspace_id,plan_id,task_id,seq,event_type,payload,created_by,created_at
            FROM work_event WHERE plan_id=%s AND workspace_id=%s AND seq>%s
            ORDER BY seq LIMIT %s
            """,
            (plan_id, actor.workspace_id, after, limit),
        )
        rows = await result.fetchall()
    # Contract《720》listPlanEvents: ListQuery -> ListResponse<WorkEventDTO>
    data = [event_view(row) for row in rows]
    return {
        "items": data,
        "data": data,
        "next_cursor": None,
        "has_more": False,
        "meta": {"request_id": None, "count": len(data)},
    }


def _sse_event(row: dict[str, Any]) -> str:
    created_at = row.get("created_at")
    if hasattr(created_at, "isoformat"):
        created_at = created_at.isoformat()
    payload = {
        "id": row["id"],
        "workspace_id": row["workspace_id"],
        "plan_id": row["plan_id"],
        "task_id": row.get("task_id"),
        "seq": row["seq"],
        "event_type": row["event_type"],
        "payload": redact_sensitive(row.get("payload") or {}),
        "created_by": row["created_by"],
        "created_at": created_at,
    }
    return f"id: {row['seq']}\nevent: {row['event_type']}\ndata: {json_dumps(payload)}\n\n"


@router.get("/plans/{plan_id}/events/stream")
async def stream_plan_events(
    plan_id: str,
    actor: Annotated[Actor, Depends(get_actor)],
    after: int = Query(default=0, ge=0),
    timeout_seconds: int = Query(default=60, ge=1, le=120),
):
    _require(actor, "read")
    async with pool.connection() as conn:
        await _owned_plan(conn, plan_id, actor)

    async def event_stream():
        cursor = after
        deadline = monotonic() + timeout_seconds
        terminal_types = {
            "plan.execution.completed",
            "plan.execution.failed",
            "plan.execution.cancelled",
        }
        while monotonic() < deadline:
            async with pool.connection() as conn:
                result = await conn.execute(
                    """
                    SELECT id,workspace_id,plan_id,task_id,seq,event_type,payload,created_by,created_at
                    FROM work_event
                    WHERE plan_id=%s AND workspace_id=%s AND seq>%s
                    ORDER BY seq LIMIT 100
                    """,
                    (plan_id, actor.workspace_id, cursor),
                )
                rows = await result.fetchall()
            if rows:
                for row in rows:
                    cursor = row["seq"]
                    yield _sse_event(row)
                    if row["event_type"] in terminal_types:
                        return
                continue
            yield ": workama-heartbeat\n\n"
            await asyncio.sleep(0.5)
        yield "event: work.stream.timeout\ndata: {\"plan_id\": " + json_dumps(plan_id) + "}\n\n"

    return StreamingResponse(
        event_stream(),
        media_type="text/event-stream",
        headers={"Cache-Control": "no-cache", "Connection": "keep-alive", "X-Accel-Buffering": "no"},
    )


@router.post("/plans/{plan_id}/sources", status_code=status.HTTP_201_CREATED)
async def add_source(plan_id: str, body: ResearchSourceCreate, actor: Annotated[Actor, Depends(get_actor)]):
    _require(actor, "write")
    validated = validate_research_url(body.url, resolve_dns=not body.url.strip().lower().startswith("mock://"))
    if body.fetch:
        # https 源走 sandbox-browser 真实抓取；mock 源保留确定性 fixture
        if validated["source_type"] == "https":
            fetched = await sandbox_browser_fetch(validated["url"])
        else:
            fetched = deterministic_mock_browser_fetch(validated["url"])
    else:
        fetched = None
    title = body.title or (fetched["title"] if fetched else None)
    excerpt = body.excerpt or (fetched["text"] if fetched else None)
    content_sha256 = fetched["content_sha256"] if fetched else None
    async with pool.connection() as conn:
        async with conn.transaction():
            plan = await _owned_plan(conn, plan_id, actor, for_update=True)
            result = await conn.execute(
                """
                INSERT INTO work_citation(id,workspace_id,plan_id,source_type,url,title,excerpt,content_sha256,created_by)
                VALUES (%s,%s,%s,%s,%s,%s,%s,%s,%s)
                ON CONFLICT(plan_id,url) DO NOTHING
                RETURNING id,plan_id,task_id,source_type,url,title,excerpt,content_sha256,created_by,created_at
                """,
                (new_id("wsrc"), actor.workspace_id, plan_id, validated["source_type"], validated["url"], title, excerpt, content_sha256, actor.user_id),
            )
            source = await result.fetchone()
            if not source:
                existing = await conn.execute(
                    "SELECT id,plan_id,task_id,source_type,url,title,excerpt,content_sha256,created_by,created_at FROM work_citation WHERE plan_id=%s AND url=%s",
                    (plan_id, validated["url"]),
                )
                source = await existing.fetchone()
            event = await _append_event(conn, plan=plan, actor=actor, event_type="citation.created", payload={"citation_id": source["id"], "source_type": validated["source_type"], "url": validated["url"]})
    return {**source, "event": event, "fetched": fetched}


@router.get("/plans/{plan_id}/sources")
async def list_sources(plan_id: str, actor: Annotated[Actor, Depends(get_actor)]):
    _require(actor, "read")
    async with pool.connection() as conn:
        await _owned_plan(conn, plan_id, actor)
        result = await conn.execute(
            "SELECT id,plan_id,task_id,source_type,url,title,excerpt,content_sha256,created_by,created_at FROM work_citation WHERE plan_id=%s AND workspace_id=%s ORDER BY created_at,id",
            (plan_id, actor.workspace_id),
        )
        rows = await result.fetchall()
    # Contract《720》listPlanSources: ListQuery -> ListResponse<WorkCitationDTO>
    data = list(rows)
    return {
        "items": data,
        "data": data,
        "next_cursor": None,
        "has_more": False,
        "meta": {"request_id": None, "count": len(data)},
    }


@router.post("/plans/{plan_id}/artifacts", status_code=status.HTTP_201_CREATED)
async def create_office_artifact(
    plan_id: str,
    body: OfficeArtifactCreate,
    actor: Annotated[Actor, Depends(get_actor)],
):
    _require(actor, "write")
    async with pool.connection() as conn:
        plan = await _owned_plan(conn, plan_id, actor)
        tasks_result = await conn.execute(
            "SELECT id,title,description,position,status FROM work_task WHERE plan_id=%s AND workspace_id=%s ORDER BY position,id",
            (plan_id, actor.workspace_id),
        )
        sources_result = await conn.execute(
            "SELECT id,url,title,excerpt,content_sha256 FROM work_citation WHERE plan_id=%s AND workspace_id=%s ORDER BY created_at,id",
            (plan_id, actor.workspace_id),
        )
        tasks = await tasks_result.fetchall()
        sources = await sources_result.fetchall()
    generated = generate_office_artifact(body.format, plan["title"], plan["objective"], tasks, sources)
    work_artifact_id = new_id("wart")
    filename = _safe_filename(body.filename or plan["title"], generated.extension)
    content_sha256 = hashlib.sha256(generated.data).hexdigest()
    s3_key = f"artifacts/{actor.workspace_id}/{work_artifact_id}/v1/{filename}"
    if body.upload:
        await put_object(ARTIFACT_BUCKET, s3_key, generated.data)
    preview = {"format": body.format, "task_count": len(tasks), "source_count": len(sources), "content_sha256": content_sha256}
    async with pool.connection() as conn:
        async with conn.transaction():
            plan = await _owned_plan(conn, plan_id, actor, for_update=True)
            artifact_id: str | None = None
            if plan.get("session_id"):
                artifact_id = new_id("art")
                await conn.execute(
                    """
                    INSERT INTO ag_artifact(
                        id,session_id,workspace_id,name,content_type,content,kind,s3_key,
                        size_bytes,content_sha256,status,preview,created_at
                    ) VALUES (%s,%s,%s,%s,%s,'','office',%s,%s,%s,%s,%s::jsonb,now())
                    """,
                    (artifact_id, plan["session_id"], actor.workspace_id, filename, generated.content_type, s3_key, len(generated.data), content_sha256, "ready" if body.upload else "pending", json_dumps(preview)),
                )
            await conn.execute(
                """
                INSERT INTO work_artifact(
                    id,workspace_id,plan_id,artifact_id,name,kind,content_type,s3_key,
                    size_bytes,content_sha256,status,preview,created_by
                ) VALUES (%s,%s,%s,%s,%s,'office',%s,%s,%s,%s,%s,%s::jsonb,%s)
                """,
                (work_artifact_id, actor.workspace_id, plan_id, artifact_id, filename, generated.content_type, s3_key, len(generated.data), content_sha256, "ready" if body.upload else "pending", json_dumps(preview), actor.user_id),
            )
            event = await _append_event(conn, plan=plan, actor=actor, event_type="artifact.created", payload={"artifact_id": work_artifact_id, "format": body.format, "name": filename, "content_sha256": content_sha256})
    return {
        "id": work_artifact_id,
        "artifact_id": artifact_id,
        "plan_id": plan_id,
        "name": filename,
        "kind": "office",
        "content_type": generated.content_type,
        "s3_key": s3_key,
        "size_bytes": len(generated.data),
        "content_sha256": content_sha256,
        "status": "ready" if body.upload else "pending",
        "preview": preview,
        "event": event,
    }


@router.get("/plans/{plan_id}/artifacts")
async def list_work_artifacts(plan_id: str, actor: Annotated[Actor, Depends(get_actor)]):
    _require(actor, "read")
    async with pool.connection() as conn:
        await _owned_plan(conn, plan_id, actor)
        result = await conn.execute(
            """
            SELECT id,workspace_id,plan_id,task_id,artifact_id,name,kind,content_type,
                   s3_key,size_bytes,content_sha256,status,preview,created_by,created_at
            FROM work_artifact
            WHERE plan_id=%s AND workspace_id=%s
            ORDER BY created_at DESC, id DESC
            """,
            (plan_id, actor.workspace_id),
        )
        rows = await result.fetchall()
    # Contract《720》listPlanArtifacts: ListQuery -> ListResponse<WorkArtifactDTO>
    data = list(rows)
    return {
        "items": data,
        "data": data,
        "next_cursor": None,
        "has_more": False,
        "meta": {"request_id": None, "count": len(data)},
    }


@router.get("/plans/{plan_id}/artifacts/{work_artifact_id}/content")
async def download_work_artifact(
    plan_id: str,
    work_artifact_id: str,
    actor: Annotated[Actor, Depends(get_actor)],
):
    _require(actor, "read")
    async with pool.connection() as conn:
        await _owned_plan(conn, plan_id, actor)
        result = await conn.execute(
            """
            SELECT name,content_type,s3_key,status
            FROM work_artifact
            WHERE id=%s AND plan_id=%s AND workspace_id=%s
            """,
            (work_artifact_id, plan_id, actor.workspace_id),
        )
        row = await result.fetchone()
    if not row:
        raise HTTPException(status_code=404, detail="Work artifact not found")
    if row["status"] != "ready":
        raise HTTPException(status_code=409, detail="Work artifact is not ready")
    try:
        content = await get_object(ARTIFACT_BUCKET, row["s3_key"])
    except Exception as exc:
        raise HTTPException(status_code=503, detail="Work artifact content is unavailable") from exc
    filename = re.sub(r"[^A-Za-z0-9._-]+", "-", str(row["name"])).strip(".-") or work_artifact_id
    return Response(
        content=content,
        media_type=row["content_type"],
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )
